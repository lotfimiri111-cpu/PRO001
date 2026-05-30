"""
core/preview.py v6 — نظام معاينة احترافي بجودة حقيقية
يستخدم LibreOffice Headless لتحويل PPTX → PDF → صور عالية الجودة
مع حماية كاملة للملف الأصلي
"""
import base64
import hashlib
import hmac
import io
import logging
import os
import secrets
import shutil
import subprocess
import tempfile
import threading
import time
from pathlib import Path
from typing import Dict, List, Optional, Tuple

log = logging.getLogger(__name__)

# ─── Storage & Security ───────────────────────────────────────────────────────
_preview_store: Dict[str, dict] = {}
_store_lock = threading.Lock()

PREVIEW_TTL = 3 * 3600      # صور المعاينة تُحذف بعد 3 ساعات
MAX_PREVIEW_SLIDES = 999    # جميع الشرائح
SLIDE_QUALITY_DPI = 150     # DPI للصور (جودة عالية مع حجم معقول)
WATERMARK_TEXT = "مذكرتي Pro — معاينة فقط"

# ─── Session Management ───────────────────────────────────────────────────────

def create_preview_session(presentation_id: str) -> str:
    """ينشئ session_id مؤقت للمعاينة ويرجع token آمن"""
    token = secrets.token_urlsafe(32)
    with _store_lock:
        _preview_store[presentation_id] = {
            "token": token,
            "slides": [],           # base64 صور
            "status": "pending",    # pending | ready | error
            "created_at": time.time(),
            "slide_count": 0,
        }
    return token


def get_preview_session(presentation_id: str) -> Optional[dict]:
    _cleanup_expired()
    with _store_lock:
        return _preview_store.get(presentation_id)


def get_preview_slides(presentation_id: str, token: str) -> Optional[List[str]]:
    """يرجع الصور فقط إذا كان التوكن صحيحاً — signed URL simulation"""
    session = get_preview_session(presentation_id)
    if not session:
        return None
    if not hmac.compare_digest(session["token"], token):
        log.warning(f"Invalid token for preview {presentation_id}")
        return None
    return session.get("slides", [])


def set_preview_ready(presentation_id: str, slides: List[str]):
    with _store_lock:
        if presentation_id in _preview_store:
            _preview_store[presentation_id]["slides"] = slides
            _preview_store[presentation_id]["status"] = "ready"
            _preview_store[presentation_id]["slide_count"] = len(slides)


def set_preview_error(presentation_id: str, msg: str):
    with _store_lock:
        if presentation_id in _preview_store:
            _preview_store[presentation_id]["status"] = "error"
            _preview_store[presentation_id]["error"] = msg


def _cleanup_expired():
    now = time.time()
    with _store_lock:
        expired = [k for k, v in _preview_store.items()
                   if now - v.get("created_at", 0) > PREVIEW_TTL]
        for k in expired:
            del _preview_store[k]

# ─── LibreOffice Conversion ───────────────────────────────────────────────────

def _find_soffice() -> Optional[str]:
    """يجد مسار LibreOffice"""
    candidates = [
        "/opt/render/project/src/scripts/office/soffice.py",
        os.path.join(os.path.dirname(os.path.dirname(__file__)), "scripts", "office", "soffice.py"),
        "soffice",
    ]
    # Python wrapper
    for c in candidates[:2]:
        if os.path.exists(c):
            return ("python", c)
    # System soffice
    result = shutil.which("soffice") or shutil.which("libreoffice")
    if result:
        return result
    return None


def _pptx_to_pdf_via_libreoffice(pptx_path: str, output_dir: str) -> Optional[str]:
    """يحوّل PPTX → PDF باستخدام LibreOffice"""
    soffice = _find_soffice()
    if not soffice:
        log.warning("LibreOffice not found, falling back to Pillow renderer")
        return None

    try:
        if isinstance(soffice, tuple):
            cmd = list(soffice) + ["--headless", "--convert-to", "pdf",
                                    "--outdir", output_dir, pptx_path]
        else:
            cmd = [soffice, "--headless", "--convert-to", "pdf",
                   "--outdir", output_dir, pptx_path]

        result = subprocess.run(
            cmd, capture_output=True, text=True, timeout=120
        )
        log.info(f"soffice stdout: {result.stdout[:200]}")
        log.info(f"soffice stderr: {result.stderr[:200]}")

        pdf_name = Path(pptx_path).stem + ".pdf"
        pdf_path = os.path.join(output_dir, pdf_name)
        if os.path.exists(pdf_path):
            return pdf_path
        # Try to find any PDF in output dir
        pdfs = list(Path(output_dir).glob("*.pdf"))
        if pdfs:
            return str(pdfs[0])
    except Exception as e:
        log.warning(f"LibreOffice conversion failed: {e}")
    return None


def _pdf_to_images(pdf_path: str, output_dir: str, dpi: int = 150) -> List[str]:
    """يحوّل PDF → صور PNG باستخدام pdftoppm"""
    prefix = os.path.join(output_dir, "slide")
    try:
        result = subprocess.run(
            ["pdftoppm", "-png", "-r", str(dpi), pdf_path, prefix],
            capture_output=True, text=True, timeout=120
        )
        if result.returncode == 0:
            images = sorted(Path(output_dir).glob("slide-*.png")) or \
                     sorted(Path(output_dir).glob("slide*.png"))
            return [str(p) for p in images]
    except Exception as e:
        log.warning(f"pdftoppm failed: {e}")
    return []


def _add_watermark_to_image(img_path: str) -> Optional[bytes]:
    """يضيف watermark شفاف متكرر للصورة"""
    try:
        from PIL import Image, ImageDraw, ImageFont
        img = Image.open(img_path).convert("RGBA")
        W, H = img.size

        overlay = Image.new("RGBA", (W, H), (0, 0, 0, 0))
        draw = ImageDraw.Draw(overlay)

        # تحديد حجم الخط
        font_size = max(20, W // 30)
        font = _find_font(font_size, bold=True)
        text = WATERMARK_TEXT

        try:
            bb = draw.textbbox((0, 0), text, font=font)
            tw, th = bb[2] - bb[0], bb[3] - bb[1]
        except:
            tw, th = len(text) * (font_size // 2), font_size + 4

        # رسم watermark مائل متكرر
        pad = 40
        tl = Image.new("RGBA", (tw + pad * 2, th + pad), (0, 0, 0, 0))
        td = ImageDraw.Draw(tl)
        td.text((pad, pad // 2), text, fill=(255, 255, 255, 70), font=font)
        rot = tl.rotate(-30, expand=True)
        rw, rh = rot.size

        step_x = max(rw + 60, W // 3)
        step_y = max(rh + 40, H // 4)

        for row in range(-1, H // step_y + 2):
            for col in range(-1, W // step_x + 2):
                x = col * step_x - rw // 4
                y = row * step_y - rh // 4
                overlay.paste(rot, (x, y), rot)

        result = Image.alpha_composite(img, overlay).convert("RGB")
        buf = io.BytesIO()
        result.save(buf, "WEBP", quality=88, method=4)
        return buf.getvalue()

    except Exception as e:
        log.warning(f"Watermark failed for {img_path}: {e}")
        # Return image without watermark as fallback
        try:
            with open(img_path, "rb") as f:
                return f.read()
        except:
            return None


# ─── Fallback Pillow Renderer ─────────────────────────────────────────────────
# (استخدامه كـ fallback فقط إذا لم يكن LibreOffice متاحاً)

CAIRO_PATHS = [
    os.path.expanduser("~/.fonts/cairo/Cairo.ttf"),
    "/root/.fonts/cairo/Cairo.ttf",
    "/tmp/fonts/cairo/Cairo.ttf",
    "/opt/render/project/src/.fonts/cairo/Cairo.ttf",
]
FALLBACK_FONTS = [
    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
]


def _find_font(size=16, bold=False):
    from PIL import ImageFont
    for p in CAIRO_PATHS:
        if os.path.exists(p):
            try:
                return ImageFont.truetype(p, size)
            except:
                pass
    for p in FALLBACK_FONTS:
        if os.path.exists(p):
            try:
                return ImageFont.truetype(p, size)
            except:
                pass
    return ImageFont.load_default()


def _pillow_fallback_render(pptx_path: str) -> List[bytes]:
    """Fallback: رسم الشرائح بـ Pillow إذا فشل LibreOffice"""
    try:
        from pptx import Presentation
        from PIL import Image, ImageDraw

        prs = Presentation(pptx_path)
        W = max(1280, int((prs.slide_width or 9144000) / 9525))
        H = max(720, int((prs.slide_height or 5143500) / 9525))
        results = []

        for slide in prs.slides:
            img = _render_slide_pillow(slide, W, H)
            wm_bytes = _watermark_pil_img(img)
            results.append(wm_bytes)
        return results
    except Exception as e:
        log.error(f"Pillow fallback failed: {e}")
        return []


def _render_slide_pillow(slide, W, H):
    """رسم شريحة واحدة بـ Pillow"""
    from PIL import Image, ImageDraw
    bg = _bg_color(slide)
    img = _draw_bg(bg, W, H, slide)
    draw = ImageDraw.Draw(img)
    for shape in slide.shapes:
        try:
            _draw_shape(draw, img, shape, W, H, bg)
        except:
            pass
    return img


def _watermark_pil_img(img) -> bytes:
    """يضيف watermark لصورة PIL"""
    from PIL import Image, ImageDraw
    W, H = img.size
    overlay = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    d = ImageDraw.Draw(overlay)
    font = _find_font(max(22, W // 28), bold=True)
    text = WATERMARK_TEXT
    try:
        bb = d.textbbox((0, 0), text, font=font)
        tw, th = bb[2] - bb[0], bb[3] - bb[1]
    except:
        tw, th = len(text) * 14, 28
    pad = 40
    tl = Image.new("RGBA", (tw + pad, th + pad), (0, 0, 0, 0))
    td = ImageDraw.Draw(tl)
    td.text((pad // 2, pad // 2), text, fill=(255, 255, 255, 80), font=font)
    rot = tl.rotate(-30, expand=True)
    rw, rh = rot.size
    sx = max(rw + 30, W // 3)
    sy = max(rh + 20, H // 3)
    for row in range(-1, H // sy + 2):
        for col in range(-1, W // sx + 2):
            overlay.paste(rot, (col * sx - rw // 2, row * sy - rh // 2), rot)
    out = Image.alpha_composite(img.convert("RGBA"), overlay).convert("RGB")
    buf = io.BytesIO()
    out.save(buf, "WEBP", quality=88)
    return buf.getvalue()


# ─── Color helpers (unchanged from v5) ───────────────────────────────────────

def _hex(h) -> tuple:
    h = str(h).lstrip("#")
    if len(h) == 6:
        return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    return (200, 200, 200)


def _read_color_xml(elem) -> Optional[tuple]:
    if elem is None:
        return None
    try:
        for child in elem.iter():
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag == 'srgbClr':
                val = child.get('val', '')
                if val:
                    return _hex(val)
            elif tag == 'sysClr':
                val = child.get('lastClr', '')
                if val:
                    return _hex(val)
            elif tag == 'prstClr':
                prstMap = {'white': (255, 255, 255), 'black': (0, 0, 0),
                           'red': (255, 0, 0), 'blue': (0, 0, 255)}
                return prstMap.get(child.get('val', ''), (128, 128, 128))
    except:
        pass
    return None


def _bg_color(slide) -> tuple:
    try:
        fill = slide.background.fill
        if fill.type is not None:
            c = fill.fore_color.rgb
            return (c.r, c.g, c.b)
    except:
        pass
    try:
        bg_elem = slide.background._element
        c = _read_color_xml(bg_elem)
        if c:
            return c
    except:
        pass
    return (255, 255, 255)


def _shape_fill_color(shape) -> Optional[tuple]:
    try:
        fill = shape.fill
        from pptx.enum.dml import MSO_FILL
        if hasattr(MSO_FILL, 'BACKGROUND') and fill.type == MSO_FILL.BACKGROUND:
            return None
        if fill.type is not None and fill.type != 5:
            try:
                c = fill.fore_color.rgb
                return (c.r, c.g, c.b)
            except:
                pass
    except:
        pass
    try:
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        sp_elem = shape._element
        solid = sp_elem.find(f'.//{{{ns}}}solidFill')
        if solid is not None:
            c = _read_color_xml(solid)
            if c:
                return c
    except:
        pass
    return None


def _text_color_from_xml(run) -> Optional[tuple]:
    try:
        rgb = run.font.color.rgb
        return (rgb.r, rgb.g, rgb.b)
    except:
        pass
    try:
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        r_elem = run._r
        rPr = r_elem.find(f'{{{ns}}}rPr')
        if rPr is not None:
            solid = rPr.find(f'{{{ns}}}solidFill')
            if solid is not None:
                c = _read_color_xml(solid)
                if c:
                    return c
    except:
        pass
    return None


def _smart_text_color(bg: tuple) -> tuple:
    lum = 0.299 * bg[0] + 0.587 * bg[1] + 0.114 * bg[2]
    return (255, 255, 255) if lum < 128 else (20, 20, 20)


def _draw_bg(bg_color, W, H, slide):
    from PIL import Image
    try:
        from lxml import etree
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        bg_elem = slide.background._element
        grad = bg_elem.find(f'.//{{{ns}}}gradFill')
        if grad is not None:
            stops = []
            for gs in grad.findall(f'.//{{{ns}}}gs'):
                pos = int(gs.get('pos', '0')) / 100000.0
                c = _read_color_xml(gs)
                if c:
                    stops.append((pos, c))
            if len(stops) >= 2:
                return _gradient_img(W, H, sorted(stops, key=lambda x: x[0]))
    except:
        pass
    img = Image.new("RGB", (W, H), bg_color)
    return img


def _gradient_img(W, H, stops):
    from PIL import Image
    img = Image.new("RGB", (W, H))
    px = img.load()
    for y in range(H):
        t = y / H
        c1 = stops[0][1]
        for i in range(len(stops) - 1):
            p0, col0 = stops[i]
            p1, col1 = stops[i + 1]
            if p0 <= t <= p1:
                lt = (t - p0) / (p1 - p0) if p1 > p0 else 0
                c1 = tuple(int(col0[j] + (col1[j] - col0[j]) * lt) for j in range(3))
                break
        for x in range(W):
            px[x, y] = c1
    return img


def _draw_shape(draw, img, shape, W, H, slide_bg):
    EMU = 914400.0
    DPI = 96
    L = int((shape.left or 0) / EMU * DPI)
    T = int((shape.top or 0) / EMU * DPI)
    SW = int((shape.width or 0) / EMU * DPI)
    SH = int((shape.height or 0) / EMU * DPI)

    fc = _shape_fill_color(shape)
    if fc and SW > 0 and SH > 0:
        draw.rectangle([L, T, L + SW, T + SH], fill=fc)

    if not shape.has_text_frame:
        return

    bg_for_contrast = fc if fc else slide_bg
    default_color = _smart_text_color(bg_for_contrast)

    y = T + 6
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            y += 8
            continue
        fs, tc, bold = 14, None, False
        for run in para.runs:
            try:
                if run.font.size:
                    fs = max(8, min(int(run.font.size / 12700), 60))
                if run.font.bold:
                    bold = True
                c = _text_color_from_xml(run)
                if c:
                    tc = c
            except:
                pass
            break
        if tc is None:
            tc = default_color
        font = _find_font(fs, bold)
        lines = _wrap(text, font, SW - 12)
        for line in lines:
            if y >= T + SH + 10:
                break
            draw.text((L + 6, y), line, fill=tc, font=font)
            y += fs + 4
        if y >= T + SH + 10:
            break


def _wrap(text, font, max_w):
    from PIL import ImageDraw, Image
    if max_w <= 20:
        return [text]
    d = ImageDraw.Draw(Image.new("RGB", (1, 1)))

    def tw(t):
        try:
            bb = d.textbbox((0, 0), t, font=font)
            return bb[2] - bb[0]
        except:
            return len(t) * 8

    if tw(text) <= max_w:
        return [text]
    words = text.split()
    lines, cur = [], ""
    for w in words:
        test = (cur + " " + w).strip()
        if tw(test) <= max_w:
            cur = test
        else:
            if cur:
                lines.append(cur)
            cur = w
    if cur:
        lines.append(cur)
    return lines or [text]


# ─── Main Public API ──────────────────────────────────────────────────────────

def generate_preview_async(presentation_id: str, pptx_path: str) -> str:
    """
    يبدأ توليد المعاينة في الخلفية ويرجع token فوراً.
    الواجهة الأمامية تستخدم هذا التوكن للحصول على الصور بعدها.
    """
    token = create_preview_session(presentation_id)

    def _worker():
        try:
            slides_b64 = _generate_slides(pptx_path)
            set_preview_ready(presentation_id, slides_b64)
            log.info(f"Preview ready: {presentation_id} ({len(slides_b64)} slides)")
        except Exception as e:
            set_preview_error(presentation_id, str(e))
            log.error(f"Preview generation failed for {presentation_id}: {e}")

    t = threading.Thread(target=_worker, daemon=True)
    t.start()
    return token


def generate_preview_sync(presentation_id: str, pptx_path: str) -> Tuple[str, List[str]]:
    """
    يولّد المعاينة بشكل متزامن ويرجع (token, slides_b64).
    للاستخدام عند الحاجة للنتائج الفورية.
    """
    token = create_preview_session(presentation_id)
    try:
        slides_b64 = _generate_slides(pptx_path)
        set_preview_ready(presentation_id, slides_b64)
        log.info(f"Preview ready (sync): {presentation_id} ({len(slides_b64)} slides)")
    except Exception as e:
        set_preview_error(presentation_id, str(e))
        log.error(f"Preview sync failed for {presentation_id}: {e}")
        slides_b64 = []
    return token, slides_b64


def _generate_slides(pptx_path: str) -> List[str]:
    """يحوّل PPTX → صور base64 (WebP) مع watermark"""
    with tempfile.TemporaryDirectory() as tmpdir:
        # محاولة 1: LibreOffice → PDF → Images (أفضل جودة)
        pdf_path = _pptx_to_pdf_via_libreoffice(pptx_path, tmpdir)
        if pdf_path:
            log.info(f"Using LibreOffice PDF path: {pdf_path}")
            image_paths = _pdf_to_images(pdf_path, tmpdir, dpi=SLIDE_QUALITY_DPI)
            if image_paths:
                log.info(f"Generated {len(image_paths)} slide images via LibreOffice")
                slides_b64 = []
                for img_path in image_paths:
                    wm_bytes = _add_watermark_to_image(img_path)
                    if wm_bytes:
                        slides_b64.append(base64.b64encode(wm_bytes).decode())
                if slides_b64:
                    return slides_b64

        # محاولة 2: Pillow Fallback
        log.info("Falling back to Pillow renderer")
        slide_bytes_list = _pillow_fallback_render(pptx_path)
        return [base64.b64encode(b).decode() for b in slide_bytes_list]


# ─── Legacy compatibility (للكود القديم) ─────────────────────────────────────

def pptx_to_preview_images(pptx_path: str, watermark: bool = True) -> List[str]:
    """واجهة متوافقة مع الكود القديم — تُرجع base64 JPEG"""
    _, slides = generate_preview_sync("_legacy_" + str(time.time()), pptx_path)
    return slides


def get_cached_preview(pid):
    session = get_preview_session(pid)
    if session and session.get("status") == "ready":
        return session.get("slides", [])
    return None


def set_cached_preview(pid, slides):
    with _store_lock:
        if pid not in _preview_store:
            _preview_store[pid] = {
                "token": secrets.token_urlsafe(32),
                "created_at": time.time(),
            }
        _preview_store[pid]["slides"] = slides
        _preview_store[pid]["status"] = "ready"
        _preview_store[pid]["slide_count"] = len(slides)
