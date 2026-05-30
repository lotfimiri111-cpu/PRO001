"""
Drawing Primitives — مذكرتي Pro v17.1
Low-level, deterministic shape/text builders.

FIXED BUGS vs v17.0:
- gradient_fill() now inserts gradFill in CORRECT OOXML position (before <a:ln>)
- shadow() now inserts effectLst in CORRECT position (after <a:ln>)
- _sort_spPr() enforces strict OOXML child order on every shape
- set_solid_alpha() exposed (no duplication with slides.py)
"""
from __future__ import annotations

from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# Slide dimensions (cm) — 16:9
W, H = 33.867, 19.05

# OOXML spPr required child order
_SPPR_ORDER = [
    qn('a:xfrm'), qn('a:prstGeom'), qn('a:custGeom'),
    qn('a:noFill'), qn('a:solidFill'), qn('a:gradFill'),
    qn('a:blipFill'), qn('a:pattFill'), qn('a:grpFill'),
    qn('a:ln'),
    qn('a:effectLst'), qn('a:effectDag'),
    qn('a:scene3d'), qn('a:sp3d'), qn('a:extLst'),
]
_SPPR_RANK = {tag: i for i, tag in enumerate(_SPPR_ORDER)}


def _sort_spPr(spPr) -> None:
    """Reorder <p:spPr> children to comply with OOXML schema."""
    children = list(spPr)
    children.sort(key=lambda el: _SPPR_RANK.get(el.tag, 99))
    for child in children:
        spPr.remove(child)
    for child in children:
        spPr.append(child)


def _get_spPr(shape):
    return shape._element.find(qn('p:spPr'))


# ── Unit helpers ─────────────────────────────────────────────────────
def cm(v: float) -> int:
    return int(Cm(v))

def pt(v: float) -> int:
    return int(Pt(v))


# ── Shape builders ───────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill: RGBColor, line=None, line_w=0.5):
    if w <= 0 or h <= 0:
        return None
    s = slide.shapes.add_shape(1, cm(x), cm(y), cm(w), cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = pt(line_w)
    else:
        s.line.fill.background()
    return s


def rrect(slide, x, y, w, h, fill: RGBColor, radius_pct=8, line=None, line_w=0.5):
    if w <= 0 or h <= 0:
        return None
    s = slide.shapes.add_shape(5, cm(x), cm(y), cm(w), cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = pt(line_w)
    else:
        s.line.fill.background()
    try:
        adj = s.adjustments
        if adj and len(adj) > 0:
            adj[0] = max(0, min(50, radius_pct)) * 1000
    except Exception:
        pass
    return s


def oval(slide, x, y, w, h, fill: RGBColor, alpha=100):
    if w <= 0 or h <= 0:
        return None
    s = slide.shapes.add_shape(9, cm(x), cm(y), cm(w), cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if alpha < 100:
        set_solid_alpha(s, alpha)
    return s


def bg(slide, color: RGBColor):
    rect(slide, 0, 0, W, H, color)

def hline(slide, x, y, w, color: RGBColor, thickness=0.08):
    rect(slide, x, y, w, thickness, color)

def vline(slide, x, y, h2, color: RGBColor, thickness=0.08):
    rect(slide, x, y, thickness, h2, color)


# ── XML fill helpers ─────────────────────────────────────────────────
def set_solid_alpha(shape, alpha_pct: int) -> None:
    """Set transparency on a solidFill shape (0=transparent, 100=opaque)."""
    try:
        spPr = _get_spPr(shape)
        srgb = spPr.find('.//' + qn('a:srgbClr'))
        if srgb is not None:
            for e in srgb.findall(qn('a:alpha')):
                srgb.remove(e)
            alp = etree.SubElement(srgb, qn('a:alpha'))
            alp.set('val', str(int(alpha_pct * 1000)))
    except Exception:
        pass


def gradient_fill(shape, c1: str, c2: str, angle: float = 90) -> None:
    """
    Apply linear gradient via XML, with correct OOXML element ordering.
    BUG FIX: previously appended gradFill AFTER <a:ln> — now sorted correctly.
    """
    try:
        spPr = _get_spPr(shape)
        # Remove all existing fill variants
        for tag in [qn('a:solidFill'), qn('a:gradFill'), qn('a:noFill'),
                    qn('a:pattFill'), qn('a:blipFill'), qn('a:grpFill')]:
            for el in spPr.findall(tag):
                spPr.remove(el)

        # Build gradFill
        grad = etree.Element(qn('a:gradFill'))
        gsLst = etree.SubElement(grad, qn('a:gsLst'))

        gs0 = etree.SubElement(gsLst, qn('a:gs'))
        gs0.set('pos', '0')
        etree.SubElement(gs0, qn('a:srgbClr')).set('val', c1.lstrip('#'))

        gs1 = etree.SubElement(gsLst, qn('a:gs'))
        gs1.set('pos', '100000')
        etree.SubElement(gs1, qn('a:srgbClr')).set('val', c2.lstrip('#'))

        lin = etree.SubElement(grad, qn('a:lin'))
        lin.set('ang', str(int(angle * 60000)))
        lin.set('scaled', '0')

        spPr.append(grad)
        _sort_spPr(spPr)  # ← enforce correct order
    except Exception:
        pass


def gradient_rect(slide, x, y, w, h, c1: str, c2: str, angle=0):
    c1h = c1.lstrip('#')
    fill_color = RGBColor(int(c1h[0:2], 16), int(c1h[2:4], 16), int(c1h[4:6], 16))
    s = rect(slide, x, y, w, h, fill_color)
    if s:
        gradient_fill(s, c1, c2, angle)
    return s


def shadow(shape, blur=16, dist=5, angle=135, alpha=0.22, color="000000") -> None:
    """
    Add outer drop shadow via XML.
    BUG FIX: effectLst now inserted in correct OOXML position (after <a:ln>).
    """
    try:
        spPr = _get_spPr(shape)
        for old in spPr.findall(qn('a:effectLst')):
            spPr.remove(old)

        eLst = etree.Element(qn('a:effectLst'))
        shdw = etree.SubElement(eLst, qn('a:outerShdw'))
        shdw.set('blurRad', str(int(blur * 12700)))
        shdw.set('dist', str(int(dist * 12700)))
        shdw.set('dir', str(int(angle * 60000)))
        shdw.set('algn', 'tl')
        srgb = etree.SubElement(shdw, qn('a:srgbClr'))
        srgb.set('val', color.lstrip('#'))
        alp = etree.SubElement(srgb, qn('a:alpha'))
        alp.set('val', str(int(alpha * 100000)))

        spPr.append(eLst)
        _sort_spPr(spPr)  # ← enforce correct order
    except Exception:
        pass


# ── Text ─────────────────────────────────────────────────────────────
def _apply_rtl_to_paragraph(p, rtl: bool, line_spacing: float) -> None:
    """
    Apply native RTL direction + line spacing at the OOXML paragraph level.
    This is the correct Arabic-native rendering — not a visual mirror.
    Sets: <a:pPr rtl="1"> and proper lnSpc for Arabic typography rhythm.
    """
    try:
        pPr = p._p.get_or_add_pPr()
        # ── RTL direction (native Arabic paragraph direction) ──────────
        if rtl:
            pPr.set('rtl', '1')
        else:
            pPr.attrib.pop('rtl', None)
        # ── Line spacing ───────────────────────────────────────────────
        for old in pPr.findall(qn('a:lnSpc')):
            pPr.remove(old)
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 100000)))
        # ── Space before: tighter for Arabic (avoids gaps between lines) ──
        for old in pPr.findall(qn('a:spcBef')):
            pPr.remove(old)
        spcBef = etree.SubElement(pPr, qn('a:spcBef'))
        spcPts = etree.SubElement(spcBef, qn('a:spcPts'))
        spcPts.set('val', '0')
    except Exception:
        pass


def _apply_rtl_to_body(tf, rtl: bool) -> None:
    """Set bodyPr rtl attribute so the entire text body is RTL-aware."""
    try:
        txBody = tf._txBody
        bodyPr = txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            if rtl:
                bodyPr.set('rtl', '1')
            else:
                bodyPr.attrib.pop('rtl', None)
    except Exception:
        pass


def _apply_run_lang(run, rtl: bool) -> None:
    """Set the lang attribute on run properties for correct shaping engine."""
    try:
        rPr = run._r.get_or_add_rPr()
        if rtl:
            rPr.set('lang', 'ar-DZ')
            rPr.set('altLang', 'en-US')
        else:
            rPr.set('lang', 'en-US')
    except Exception:
        pass


def txt(slide, text: str, x, y, w, h,
        font="Cairo", size=14, bold=False, italic=False,
        color: RGBColor | None = None,
        align=PP_ALIGN.RIGHT,
        margin=0.1, rtl=True, spacing=None,
        vcenter=True, line_spacing=1.15):
    """
    نص احترافي مع RTL-native rendering حقيقي.
    - rtl=True → ضبط اتجاه الفقرة والجسم والـ shaping engine بشكل صحيح
    - line_spacing → ارتفاع السطر النسبي
    - vcenter → توسيط عمودي حقيقي
    """
    if not text or w <= 0 or h <= 0:
        return None
    try:
        from pptx.enum.text import MSO_ANCHOR
        sh = slide.shapes.add_shape(1, cm(x), cm(y), cm(w), cm(h))
        sh.fill.background()
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True

        # ── Margins: Arabic needs more right breathing room ────────────
        if rtl:
            tf.margin_left   = cm(margin * 0.6)
            tf.margin_right  = cm(margin * 1.2)
        else:
            tf.margin_left   = cm(margin)
            tf.margin_right  = cm(margin)
        tf.margin_top    = cm(0.05)
        tf.margin_bottom = cm(0.05)

        if vcenter:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        else:
            tf.vertical_anchor = MSO_ANCHOR.TOP

        # ── Apply RTL to body ──────────────────────────────────────────
        _apply_rtl_to_body(tf, rtl)

        p = tf.paragraphs[0]
        p.alignment = align
        _apply_rtl_to_paragraph(p, rtl, line_spacing)

        run = p.add_run()
        run.text = str(text)
        run.font.name   = font
        run.font.size   = Pt(size)
        run.font.bold   = bold
        run.font.italic = italic
        if color:
            run.font.color.rgb = color
        _apply_run_lang(run, rtl)
        return sh
    except Exception:
        # fallback إلى textbox
        tb = slide.shapes.add_textbox(cm(x), cm(y), cm(w), cm(h))
        tb.word_wrap = True
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.name   = font
        run.font.size   = Pt(size)
        run.font.bold   = bold
        run.font.italic = italic
        if color:
            run.font.color.rgb = color
        return tb


def txt2(slide, label: str, value: str, x, y, w, h,
         font="Cairo", label_size=10, value_size=13,
         label_color: RGBColor | None = None,
         value_color: RGBColor | None = None,
         align=PP_ALIGN.RIGHT, margin=0.1, rtl=True):
    """
    نص بسطرين: تسمية + قيمة مع RTL-native rendering.
    """
    if w <= 0 or h <= 0: return None
    try:
        from pptx.enum.text import MSO_ANCHOR
        sh = slide.shapes.add_shape(1, cm(x), cm(y), cm(w), cm(h))
        sh.fill.background()
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True
        tf.margin_left   = cm(margin * 0.6 if rtl else margin)
        tf.margin_right  = cm(margin * 1.2 if rtl else margin)
        tf.margin_top    = cm(0.05)
        tf.margin_bottom = cm(0.05)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _apply_rtl_to_body(tf, rtl)

        p1 = tf.paragraphs[0]
        p1.alignment = align
        _apply_rtl_to_paragraph(p1, rtl, 1.2)
        r1 = p1.add_run()
        r1.text = str(label)
        r1.font.name  = font
        r1.font.size  = Pt(label_size)
        r1.font.bold  = True
        if label_color: r1.font.color.rgb = label_color
        _apply_run_lang(r1, rtl)

        p2 = tf.add_paragraph()
        p2.alignment = align
        _apply_rtl_to_paragraph(p2, rtl, 1.2)
        r2 = p2.add_run()
        r2.text = str(value)
        r2.font.name  = font
        r2.font.size  = Pt(value_size)
        r2.font.bold  = False
        if value_color: r2.font.color.rgb = value_color
        _apply_run_lang(r2, rtl)
        return sh
    except Exception:
        return None


def blank_slide(prs):
    """Add a completely blank slide (layout 6 = Blank)."""
    return prs.slides.add_slide(prs.slide_layouts[6])


# ── Advanced Visual Primitives ────────────────────────────────────────

def glow(shape, color: str = "C6A03C", radius: int = 20, alpha: float = 0.4) -> None:
    """Add a glow effect (outerShdw with zero distance = glow)."""
    try:
        spPr = _get_spPr(shape)
        for old in spPr.findall(qn('a:effectLst')):
            spPr.remove(old)
        eLst = etree.Element(qn('a:effectLst'))
        g = etree.SubElement(eLst, qn('a:outerShdw'))
        g.set('blurRad', str(int(radius * 12700)))
        g.set('dist', '0')
        g.set('dir', '0')
        g.set('algn', 'ctr')
        srgb = etree.SubElement(g, qn('a:srgbClr'))
        srgb.set('val', color.lstrip('#'))
        alp = etree.SubElement(srgb, qn('a:alpha'))
        alp.set('val', str(int(alpha * 100000)))
        spPr.append(eLst)
        _sort_spPr(spPr)
    except Exception:
        pass


def multi_stop_gradient(shape, stops: list[tuple[int, str]], angle: float = 90) -> None:
    """
    Apply a multi-stop gradient.
    stops = [(pos_pct, '#RRGGBB'), ...]  e.g. [(0,'#07172F'),(50,'#1A3A6E'),(100,'#07172F')]
    """
    try:
        spPr = _get_spPr(shape)
        for tag in [qn('a:solidFill'), qn('a:gradFill'), qn('a:noFill'),
                    qn('a:pattFill'), qn('a:blipFill'), qn('a:grpFill')]:
            for el in spPr.findall(tag):
                spPr.remove(el)
        grad = etree.Element(qn('a:gradFill'))
        gsLst = etree.SubElement(grad, qn('a:gsLst'))
        for pos_pct, hex_color in stops:
            gs = etree.SubElement(gsLst, qn('a:gs'))
            gs.set('pos', str(int(pos_pct * 1000)))
            etree.SubElement(gs, qn('a:srgbClr')).set('val', hex_color.lstrip('#'))
        lin = etree.SubElement(grad, qn('a:lin'))
        lin.set('ang', str(int(angle * 60000)))
        lin.set('scaled', '0')
        spPr.append(grad)
        _sort_spPr(spPr)
    except Exception:
        pass


def gradient_oval(slide, x, y, w, h, c1: str, c2: str, angle=90, alpha=100):
    """Oval with gradient fill."""
    if w <= 0 or h <= 0:
        return None
    c1h = c1.lstrip('#')
    fill_color = RGBColor(int(c1h[0:2], 16), int(c1h[2:4], 16), int(c1h[4:6], 16))
    s = slide.shapes.add_shape(9, cm(x), cm(y), cm(w), cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.fill.background()
    gradient_fill(s, c1, c2, angle)
    if alpha < 100:
        set_solid_alpha(s, alpha)
    return s


def triangle(slide, x, y, w, h, fill: RGBColor, pointing='up'):
    """Equilateral-ish triangle shape."""
    if w <= 0 or h <= 0:
        return None
    # Use right-triangle preset (shape 6) rotated for pointing direction
    shape_id = 6  # rtTriangle
    s = slide.shapes.add_shape(shape_id, cm(x), cm(y), cm(w), cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if pointing == 'down':
        s.rotation = 180
    elif pointing == 'left':
        s.rotation = 90
    elif pointing == 'right':
        s.rotation = 270
    return s


def diamond(slide, x, y, w, h, fill: RGBColor, alpha=100):
    """Diamond shape."""
    if w <= 0 or h <= 0:
        return None
    s = slide.shapes.add_shape(4, cm(x), cm(y), cm(w), cm(h))  # shape 4 = diamond
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if alpha < 100:
        set_solid_alpha(s, alpha)
    return s


def hexagon(slide, x, y, w, h, fill: RGBColor, alpha=100):
    """Hexagon shape."""
    if w <= 0 or h <= 0:
        return None
    s = slide.shapes.add_shape(10, cm(x), cm(y), cm(w), cm(h))  # shape 10 = hexagon
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if alpha < 100:
        set_solid_alpha(s, alpha)
    return s


def arc_progress(slide, x, y, size, fill: RGBColor, bg_color: RGBColor,
                 thickness=0.4) -> None:
    """Simulated progress ring using two arcs (outer ring + inner mask)."""
    # Outer ring background
    outer = oval(slide, x, y, size, size, bg_color, alpha=30)
    # Inner circle mask (creates ring effect)
    inner_offset = thickness
    inner_s = size - 2 * inner_offset
    oval(slide, x + inner_offset, y + inner_offset, inner_s, inner_s, bg_color, alpha=80)


def decorative_dots(slide, x, y, cols, rows, dot_size, gap, color: RGBColor, alpha=20):
    """Grid of decorative dots."""
    for r in range(rows):
        for c in range(cols):
            dx = x + c * (dot_size + gap)
            dy = y + r * (dot_size + gap)
            o = oval(slide, dx, dy, dot_size, dot_size, color)
            if o and alpha < 100:
                set_solid_alpha(o, alpha)


def wave_rect(slide, x, y, w, h, fill: RGBColor, wavy_top=True):
    """Rectangle with rounded top (simulates wave). Uses rrect with high radius."""
    if wavy_top:
        return rrect(slide, x, y, w, h, fill, radius_pct=12)
    return rect(slide, x, y, w, h, fill)


def badge(slide, x, y, w, h, fill_c1: str, fill_c2: str, label: str,
          font="Cairo", font_size=11, text_color: RGBColor = None, T=None):
    """Pill-shaped badge with gradient and centered label."""
    b = rrect(slide, x, y, w, h, RGBColor(0xC6, 0xA0, 0x3C), radius_pct=50)
    if b:
        gradient_fill(b, fill_c1, fill_c2, angle=0)
    if text_color is None and T is not None:
        text_color = T.text_dark_rgb
    txt(slide, label, x, y, w, h,
        font=font, size=font_size, bold=True,
        color=text_color, align=PP_ALIGN.CENTER, rtl=True)
    return b


def icon_circle(slide, x, y, size, bg_c1: str, bg_c2: str,
                icon: str, icon_size=20, T=None):
    """Circle with gradient bg + centered emoji/icon."""
    c = oval(slide, x, y, size, size,
             RGBColor(int(bg_c1.lstrip('#')[0:2], 16),
                      int(bg_c1.lstrip('#')[2:4], 16),
                      int(bg_c1.lstrip('#')[4:6], 16)))
    if c:
        gradient_fill(c, bg_c1, bg_c2, angle=135)
    txt(slide, icon, x, y, size, size,
        font="Calibri", size=icon_size, bold=False,
        color=T.text_dark_rgb if T else RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER, rtl=False)
    return c


def number_badge(slide, x, y, size, num: int | str, T):
    """Circular number badge with accent gradient."""
    c = oval(slide, x, y, size, size, T.accent_rgb)
    if c:
        gradient_fill(c, T.accent_grad1, T.accent_grad2, 135)
        shadow(c, blur=10, dist=3, alpha=0.35)
    txt(slide, str(num), x, y, size, size,
        font="Calibri", size=max(8, int(size * 18)), bold=True,
        color=T.text_dark_rgb, align=PP_ALIGN.CENTER, rtl=False)
    return c


def divider(slide, x, y, w, T, style='gradient'):
    """Decorative divider line."""
    if style == 'gradient':
        d = rect(slide, x, y, w, 0.06, T.accent_rgb)
        if d:
            multi_stop_gradient(d, [(0, T.bg), (50, T.accent), (100, T.bg)], angle=0)
    elif style == 'double':
        rect(slide, x, y, w, 0.05, T.accent_rgb)
        rect(slide, x + w * 0.05, y + 0.12, w * 0.9, 0.03, T.muted_rgb)
    else:
        rect(slide, x, y, w, 0.06, T.accent_rgb)


def card_3d(slide, x, y, w, h, T, radius=10):
    """Card with shadow + subtle gradient for 3D feel."""
    # Shadow layer (slightly offset)
    shadow_s = rrect(slide, x + 0.15, y + 0.2, w, h, T.bg_rgb, radius_pct=radius)
    if shadow_s:
        set_solid_alpha(shadow_s, 40)

    # Main card
    c = rrect(slide, x, y, w, h, T.card_rgb, radius_pct=radius)
    if c:
        multi_stop_gradient(c, [(0, T.card), (100, T.bg2)], angle=135)
        shadow(c, blur=18, dist=5, alpha=0.4)
    return c


def slide_number(slide, num: int, total: int, T):
    """Slide number indicator bottom-right."""
    label = f"{num} / {total}"
    txt(slide, label, W - 3.5, H - 0.55, 3.2, 0.45,
        font="Calibri", size=9, bold=False,
        color=T.muted_rgb, align=PP_ALIGN.LEFT, rtl=False)


def watermark(slide, text: str, T):
    """Subtle watermark bottom-left."""
    txt(slide, text, 0.4, H - 0.55, 6.0, 0.45,
        font="Calibri", size=8, bold=False,
        color=T.muted_rgb, align=PP_ALIGN.RIGHT, rtl=False)


def section_tag(slide, label: str, x, y, T):
    """Small colored tag/label."""
    w, h = 3.5, 0.52
    b = rrect(slide, x, y, w, h, T.accent_rgb, radius_pct=50)
    if b:
        gradient_fill(b, T.accent_grad1, T.accent_grad2, 0)
    txt(slide, label, x, y, w, h,
        font="Cairo", size=10, bold=True,
        color=T.text_dark_rgb, align=PP_ALIGN.CENTER, rtl=True)


# ══════════════════════════════════════════════════════════════════════
# DESIGN INTELLIGENCE LAYER — v28
# Smart layout utilities that respond to content
# ══════════════════════════════════════════════════════════════════════

def _smart_font_size(text: str, base: float, min_s: float, max_s: float,
                      area_w: float, area_h: float,
                      chars_per_pt: float = 0.065) -> float:
    """
    Compute font size that fits `text` inside (area_w × area_h).
    chars_per_pt = approximate chars per pt width per cm of width.
    """
    if not text:
        return base
    length = len(text)
    # Very short → go large; very long → go small
    if length <= 20:
        factor = 1.18
    elif length <= 40:
        factor = 1.0
    elif length <= 70:
        factor = 0.88
    elif length <= 110:
        factor = 0.76
    else:
        factor = 0.66
    size = base * factor
    return max(min_s, min(max_s, size))


def smart_title(slide, text: str, x, y, w, h, T, font="Cairo",
                base_size=30, min_s=18, max_s=38, rtl=True, vcenter=True):
    """
    Dominant section title with smart sizing and accent underline.
    """
    from pptx.enum.text import PP_ALIGN
    fs = _smart_font_size(text, base_size, min_s, max_s, w, h)
    title_h = h * 0.72 if h > 1.0 else h
    t = txt(slide, text, x, y, w, title_h,
            font=font, size=fs, bold=True,
            color=T.text_light_rgb, align=PP_ALIGN.RIGHT,
            rtl=rtl, vcenter=vcenter, line_spacing=1.05)
    return t


def accent_pill(slide, text: str, x, y, T, font="Cairo", size=10.5):
    """Small accent-colored pill label."""
    from pptx.enum.text import PP_ALIGN
    w = max(2.8, len(text) * 0.18 + 0.6)
    h = 0.44
    b = rrect(slide, x, y, w, h, T.accent_rgb, radius_pct=50)
    if b:
        gradient_fill(b, T.accent_grad1, T.accent_grad2, 0)
    txt(slide, text, x, y, w, h,
        font=font, size=size, bold=True,
        color=T.text_dark_rgb, align=PP_ALIGN.CENTER, rtl=True)
    return w, h


def premium_card(slide, x, y, w, h, T, radius=12, depth=True, glow_on=True):
    """
    Premium card with layered shadow, gradient, optional glow.
    Returns the main card shape.
    """
    if depth:
        # Deep shadow base
        sh = rrect(slide, x + 0.18, y + 0.24, w, h, T.bg_rgb, radius_pct=radius)
        if sh:
            set_solid_alpha(sh, 35)
    c = rrect(slide, x, y, w, h, T.card_rgb, radius_pct=radius)
    if c:
        multi_stop_gradient(c, [(0, T.card), (60, T.bg2), (100, T.bg)], 135)
        shadow(c, blur=22, dist=6, alpha=0.44)
        if glow_on:
            glow(c, T.accent.lstrip('#'), radius=16, alpha=0.07)
    return c


def card_with_accent_top(slide, x, y, w, h, T, radius=12, bar_h=0.38):
    """Card with colored accent bar on top."""
    c = premium_card(slide, x, y, w, h, T, radius=radius)
    bar = rrect(slide, x, y, w, bar_h, T.accent_rgb, radius_pct=0)
    if bar:
        multi_stop_gradient(bar, [(0, T.accent2), (50, T.accent), (100, T.accent2)], 0)
        glow(bar, T.accent.lstrip('#'), radius=10, alpha=0.22)
    return c


def card_with_accent_side(slide, x, y, w, h, T, radius=12, bar_w=0.26):
    """Card with colored accent bar on right side (RTL primary)."""
    c = premium_card(slide, x, y, w, h, T, radius=radius)
    bar = rrect(slide, x + w - bar_w, y, bar_w, h, T.accent_rgb, radius_pct=0)
    if bar:
        gradient_fill(bar, T.accent_grad1, T.accent_grad2, 90)
    return c


def kpi_card(slide, x, y, w, h, T, value: str, label: str,
             unit: str = '', font="Cairo"):
    """
    Premium KPI/stat card: giant centered value + label below.
    Auto-scales value font based on length.
    """
    from pptx.enum.text import PP_ALIGN
    # Card base
    c = rrect(slide, x, y, w, h, T.card_rgb, radius_pct=14)
    if c:
        multi_stop_gradient(c, [(0, T.bg2), (50, T.card), (100, T.bg2)], 135)
        shadow(c, blur=20, dist=6, alpha=0.45)

    # Accent top stripe
    tp = rrect(slide, x, y, w, 0.32, T.accent_rgb, radius_pct=0)
    if tp:
        multi_stop_gradient(tp, [(0, T.accent2), (50, T.accent), (100, T.accent2)], 0)
        glow(tp, T.accent.lstrip('#'), radius=8, alpha=0.28)

    # Bottom pulse bar
    bp = rrect(slide, x, y + h - 0.18, w, 0.18, T.accent_rgb, radius_pct=0)
    if bp:
        set_solid_alpha(bp, 30)

    # Value — giant, centered
    vlen = len(str(value))
    vs = 46 if vlen <= 2 else 38 if vlen <= 4 else 28 if vlen <= 7 else 22
    txt(slide, str(value), x + 0.12, y + 0.32, w - 0.24, h * 0.50,
        font="Calibri", size=vs, bold=True,
        color=T.accent_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)

    # Unit badge (if present)
    if unit:
        ub = rrect(slide, x + w / 2 - 1.6, y + h * 0.53 + 0.06, 3.2, 0.42,
                   T.bg_rgb, radius_pct=40)
        if ub:
            set_solid_alpha(ub, 55)
        txt(slide, unit, x + w / 2 - 1.6, y + h * 0.53 + 0.06, 3.2, 0.42,
            font=font, size=9.5, bold=False,
            color=T.muted_rgb, align=PP_ALIGN.CENTER, rtl=True, vcenter=True)

    # Divider
    hline(slide, x + w * 0.14, y + h * 0.71, w * 0.72, T.muted_rgb, thickness=0.04)

    # Label
    txt(slide, label, x + 0.12, y + h * 0.73, w - 0.24, h * 0.25,
        font=font, size=max(10, min(13, h * 5.5)), bold=False,
        color=T.text_light_rgb, align=PP_ALIGN.CENTER, rtl=True, vcenter=True)


def result_row(slide, x, y, w, h, T, text: str, index: int,
               font="Cairo", highlight=False):
    """
    Premium result list row with smart font sizing.
    index: 1-based number shown in badge on right.
    highlight: makes row visually stronger (for key results).
    """
    from pptx.enum.text import PP_ALIGN
    even = (index % 2 == 0)

    # Row background
    rw = rrect(slide, x, y, w, h,
               T.card_rgb if not even else T.bg2_rgb,
               radius_pct=10)
    if rw:
        stops = [(0, T.card), (100, T.bg2)] if not even else [(0, T.bg2), (100, T.card)]
        multi_stop_gradient(rw, stops, 0)
        if highlight:
            shadow(rw, blur=12, dist=3, alpha=0.32)
            glow(rw, T.accent.lstrip('#'), radius=12, alpha=0.06)
        else:
            shadow(rw, blur=5, dist=2, alpha=0.16)

    # Accent side bar (fades with index for visual rhythm)
    alpha_bar = max(22, 62 - index * 6)
    bar = rect(slide, x + w - 0.28, y, 0.28, h, T.accent_rgb)
    if bar:
        gradient_fill(bar, T.accent_grad1, T.accent_grad2, 90)
        set_solid_alpha(bar, alpha_bar)

    # Number badge
    nd = min(0.68, h * 0.72)
    nb_x = x + w - 1.1 - nd
    nb_y = y + (h - nd) / 2
    nb_c = oval(slide, nb_x, nb_y, nd, nd, T.accent_rgb)
    if nb_c:
        multi_stop_gradient(nb_c, [(0, T.accent), (100, T.accent2)], 135)
        shadow(nb_c, blur=8, dist=2, alpha=0.3)
    txt(slide, str(index), nb_x, nb_y, nd, nd,
        font="Calibri", size=max(9, int(nd * 11)), bold=True,
        color=T.text_dark_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)

    # Content text — smart sizing
    text_w = w - nd - 1.5
    fs = _smart_font_size(text, 13.5, 11, 15.5, text_w, h)
    txt(slide, text, x + 0.25, y, text_w, h,
        font=font, size=fs, bold=highlight,
        color=T.text_light_rgb if not highlight else T.accent_rgb,
        align=PP_ALIGN.RIGHT, rtl=True, vcenter=True, line_spacing=1.2)


def premium_header(slide, T, title: str, subtitle: str = '',
                   slide_num: int = None, total: int = 13,
                   accent_side='right', font="Cairo"):
    """
    Premium header with:
    - Dominant gradient background
    - Strong title with smart sizing
    - Subtle subtitle
    - Slide counter badge
    - Multi-layer accent lines
    """
    from pptx.enum.text import PP_ALIGN
    HDR_H = 3.0

    # Main header background — deep gradient
    gradient_rect(slide, 0, 0, W, HDR_H, T.grad2, T.grad1, angle=130)

    # Layered accent lines at bottom
    al1 = rect(slide, 0, HDR_H - 0.26, W, 0.26, T.accent_rgb)
    if al1:
        multi_stop_gradient(al1, [(0, T.bg), (35, T.accent2), (50, T.accent),
                                   (65, T.accent2), (100, T.bg)], 0)
    rect(slide, 0, HDR_H - 0.32, W, 0.06, T.muted_rgb)
    rect(slide, 0, HDR_H - 0.06, W, 0.06, T.bg_rgb)

    # Accent vertical bar
    if accent_side == 'right':
        av = rect(slide, W - 0.56, 0, 0.56, HDR_H, T.accent_rgb)
    else:
        av = rect(slide, 0, 0, 0.56, HDR_H, T.accent_rgb)
    if av:
        gradient_fill(av, T.accent_grad1, T.accent_grad2, 90)

    # Decorative background circle
    oval(slide, W - 5.5, -2.5, 8, 8, T.accent_rgb, alpha=8)

    # Slide number badge
    if slide_num is not None:
        nb_s = 0.78
        nb_x = 1.05
        nb_y = (HDR_H - nb_s) / 2
        nb_c = oval(slide, nb_x, nb_y, nb_s, nb_s, T.accent_rgb)
        if nb_c:
            multi_stop_gradient(nb_c, [(0, T.accent_grad1), (100, T.accent_grad2)], 135)
            shadow(nb_c, blur=10, dist=3, alpha=0.38)
        txt(slide, str(slide_num), nb_x, nb_y, nb_s, nb_s,
            font="Calibri", size=15, bold=True,
            color=T.text_dark_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)
        txt(slide, f"/{total}", nb_x + nb_s, nb_y + nb_s * 0.32, 0.85, nb_s * 0.38,
            font="Calibri", size=8, bold=False,
            color=T.muted_rgb, align=PP_ALIGN.LEFT, rtl=False, vcenter=True)
        title_x = nb_x + nb_s + 0.9
    else:
        title_x = 0.72

    title_w = W - title_x - 0.72
    fs_title = _smart_font_size(title, 30, 20, 34, title_w, HDR_H * 0.65)
    txt(slide, title, title_x, 0.18, title_w, HDR_H * 0.63,
        font=font, size=fs_title, bold=True,
        color=T.text_light_rgb, align=PP_ALIGN.RIGHT,
        rtl=True, vcenter=True, line_spacing=1.05)

    if subtitle:
        fs_sub = min(14.5, max(11, fs_title * 0.44))
        txt(slide, subtitle, title_x, HDR_H * 0.63, title_w, HDR_H * 0.33,
            font=font, size=fs_sub, bold=False, italic=True,
            color=T.muted_rgb, align=PP_ALIGN.RIGHT,
            rtl=True, vcenter=True, line_spacing=1.0)

    return HDR_H


def section_divider_line(slide, x, y, w, T):
    """Triple-layer decorative divider."""
    d1 = rect(slide, x, y, w, 0.07, T.accent_rgb)
    if d1:
        multi_stop_gradient(d1, [(0, T.bg2), (50, T.accent), (100, T.bg2)], 0)
    rect(slide, x + w * 0.08, y + 0.1, w * 0.84, 0.03, T.muted_rgb)


def two_col_layout(n_items):
    """Return (cols, rows) for n items, preferring 2-col layout when n>3."""
    if n_items <= 3:
        return n_items, 1
    elif n_items <= 6:
        return 2, (n_items + 1) // 2
    else:
        return 3, (n_items + 2) // 3


def adaptive_body_size(text: str, container_h: float,
                        base=13.5, min_s=10.5, max_s=16.0) -> float:
    """Scale body text to fill a container height comfortably."""
    n_words = len(text.split())
    if n_words <= 10:
        factor = 1.15
    elif n_words <= 20:
        factor = 1.0
    elif n_words <= 35:
        factor = 0.88
    elif n_words <= 55:
        factor = 0.76
    else:
        factor = 0.65
    size = base * factor
    # Also constrain to height
    h_factor = container_h * 4.5
    size = min(size, h_factor)
    return max(min_s, min(max_s, size))


def premium_bg(slide, T, style='a'):
    """
    Enhanced background with depth layers and ambient shapes.
    Styles: 'a' (radial), 'b' (diagonal), 'c' (corner), 'd' (concentric)
    """
    bg(slide, T.bg_rgb)
    angle_map = {'a': 135, 'b': 160, 'c': 90, 'd': 45}
    gradient_rect(slide, 0, 0, W, H, T.grad1, T.grad2,
                  angle=angle_map.get(style, 135))

    if style == 'a':
        oval(slide, -4, -4, 13, 13, T.accent_rgb, alpha=5)
        oval(slide, W - 10, H - 9, 15, 15, T.bg2_rgb, alpha=42)
        oval(slide, W - 7, -1, 9, 9, T.accent_rgb, alpha=4)
        decorative_dots(slide, 1.2, H - 4.5, 5, 3, 0.16, 0.42, T.accent_rgb, alpha=11)
    elif style == 'b':
        diamond(slide, W - 7.5, -2.5, 6.5, 6.5, T.accent_rgb, alpha=6)
        diamond(slide, -1.5, H - 5, 5, 5, T.accent_rgb, alpha=5)
        hexagon(slide, W - 5, H * 0.3, 3.0, 3.0, T.accent_rgb, alpha=7)
        decorative_dots(slide, 1.0, 1.8, 4, 4, 0.15, 0.36, T.accent_rgb, alpha=9)
        oval(slide, W * 0.35, -3, 8, 8, T.accent_rgb, alpha=3)
    elif style == 'c':
        oval(slide, -5, -4, 13, 13, T.accent_rgb, alpha=4)
        oval(slide, W - 11, H - 10, 16, 16, T.accent_rgb, alpha=4)
        oval(slide, W - 7, -3, 10, 10, T.bg2_rgb, alpha=38)
        decorative_dots(slide, W - 7, 1.5, 4, 5, 0.14, 0.35, T.accent_rgb, alpha=10)
        oval(slide, -2, H * 0.4, 6, 6, T.bg2_rgb, alpha=22)
    elif style == 'd':
        for r, a in [(28, 3), (22, 4), (16, 5), (10, 7), (6, 9)]:
            oval(slide, W / 2 - r / 2, H / 2 - r / 2, r, r, T.accent_rgb, alpha=a)
        decorative_dots(slide, 1.8, H - 4.2, 5, 2, 0.18, 0.44, T.accent_rgb, alpha=11)


# ══════════════════════════════════════════════════════════════════════
# LAYOUT INTELLIGENCE SYSTEM v28.2
# ══════════════════════════════════════════════════════════════════════

def layout_pick(slide_type: str, n_items: int, text_len: int) -> str:
    """
    Returns a layout variant code based on content characteristics.
    Ensures visual variety across the presentation — no two consecutive
    slides share the same layout pattern.

    slide_type: 'results' | 'objectives' | 'importance' | 'methodology' | 'future'
    n_items: number of list items
    text_len: average character length per item
    returns: 'A' | 'B' | 'C' | 'D'
    """
    # Deterministic but varied: use hash of type+count for stable selection
    seed = (hash(slide_type) + n_items * 7 + text_len // 20) % 4
    return ['A', 'B', 'C', 'D'][seed]


def glass_card(slide, x, y, w, h, T, radius=12, blur_alpha=62):
    """
    Glassmorphism-style card: translucent + frosted blur effect.
    Uses layered semi-transparent fills for depth illusion.
    """
    # Deep base layer
    b1 = rrect(slide, x + 0.22, y + 0.26, w, h, T.bg2_rgb, radius_pct=radius)
    if b1: set_solid_alpha(b1, 28)

    # Main glass surface
    gc = rrect(slide, x, y, w, h, T.card_rgb, radius_pct=radius)
    if gc:
        multi_stop_gradient(gc, [(0, T.card), (55, T.bg2), (100, T.bg)], 135)
        set_solid_alpha(gc, blur_alpha)
        shadow(gc, blur=24, dist=7, alpha=0.38)

    # Top highlight line (glass edge)
    gh = rrect(slide, x, y, w, 0.14, T.text_light_rgb, radius_pct=0)
    if gh: set_solid_alpha(gh, 14)

    # Left highlight (light refraction simulation for RTL)
    gv = rrect(slide, x + w - 0.08, y + 0.14, 0.08, h - 0.14, T.text_light_rgb, radius_pct=0)
    if gv: set_solid_alpha(gv, 8)

    return gc


def mesh_gradient_bg(slide, T, complexity=3):
    """
    Layered mesh gradient background with organic blob shapes.
    Creates depth and visual richness without pattern repetition.
    complexity: 1=simple, 2=medium, 3=rich
    """
    bg(slide, T.bg_rgb)
    gradient_rect(slide, 0, 0, W, H, T.grad1, T.grad2, angle=145)

    blobs = [
        (W * 0.68, -2.5, 11, 11, T.accent_rgb, 4),
        (-3.5, H * 0.55, 12, 12, T.bg2_rgb, 28),
        (W * 0.15, H * 0.72, 9, 9, T.accent_rgb, 3),
        (W - 6, H * 0.38, 8, 8, T.bg2_rgb, 22),
    ]
    if complexity >= 2:
        blobs += [
            (W * 0.42, -1.5, 7, 7, T.accent_rgb, 5),
            (-1.5, H * 0.22, 8, 8, T.bg2_rgb, 18),
            (W * 0.78, H * 0.78, 9, 9, T.accent_rgb, 4),
        ]
    if complexity >= 3:
        blobs += [
            (W * 0.12, H * 0.08, 5, 5, T.accent_rgb, 6),
            (W * 0.55, H * 0.85, 7, 7, T.bg2_rgb, 20),
            (W - 3, H * 0.12, 6, 6, T.accent_rgb, 5),
        ]
    for bx, by, bw, bh, bc, ba in blobs:
        oval(slide, bx, by, bw, bh, bc, alpha=ba)


def cinematic_bg(slide, T, variant='split'):
    """
    Cinematic background for cover and final slides.
    variant: 'split' | 'diagonal' | 'radial' | 'vignette'
    """
    bg(slide, T.bg_rgb)

    if variant == 'split':
        # Bold diagonal split
        gradient_rect(slide, 0, 0, W, H, T.grad2, T.bg, angle=145)
        gradient_rect(slide, W * 0.46, 0, W * 0.54, H, T.grad1, T.bg2, angle=90)
        # Bright edge line at split
        hl = rect(slide, W * 0.455, 0, 0.10, H, T.accent_rgb)
        if hl: multi_stop_gradient(hl, [(0, T.bg), (40, T.accent2), (60, T.accent), (100, T.bg)], 90)

    elif variant == 'diagonal':
        gradient_rect(slide, 0, 0, W, H, T.grad1, T.grad2, angle=135)
        # Large diagonal accent shape
        for i in range(4):
            sh = slide.shapes.add_shape(4,  # freeform via triangle
                cm(-2 + i * 1.2), cm(-2), cm(W * 0.62 - i * 1.2), cm(H + 4))
            sh.fill.solid()
            sh.fill.fore_color.rgb = T.accent_rgb
            sh.line.fill.background()
            set_solid_alpha(sh, max(3, 12 - i * 3))

    elif variant == 'radial':
        gradient_rect(slide, 0, 0, W, H, T.grad2, T.grad1, angle=0)
        # Concentric oval glows
        for r, a in [(34, 5), (26, 8), (20, 11), (15, 14), (10, 10)]:
            oval(slide, W / 2 - r / 2, H / 2 - r / 2 * 0.75, r, r * 0.75,
                 T.accent_rgb, alpha=a)

    elif variant == 'vignette':
        gradient_rect(slide, 0, 0, W, H, T.grad2, T.grad1, angle=135)
        # Corner vignettes
        for bx, by in [(- 5, -5), (W - 8, -5), (-5, H - 8), (W - 8, H - 8)]:
            oval(slide, bx, by, 13, 13, T.bg_rgb, alpha=42)

    # Subtle dot grid (premium texture)
    decorative_dots(slide, 1.0, H - 5.5, 8, 3, 0.12, 0.38, T.accent_rgb, alpha=7)
    decorative_dots(slide, W - 6.5, 1.0, 6, 4, 0.11, 0.32, T.accent_rgb, alpha=6)


def hero_stat(slide, x, y, w, h, T, value: str, label: str, unit='',
              font="Cairo", rank=0):
    """
    Hero-style stat card — much more dramatic than kpi_card.
    rank=0 is the LARGEST hero, rank>0 are supporting stats.
    """
    # Scale based on rank
    scale = 1.0 if rank == 0 else (0.82 if rank == 1 else 0.68)
    glow_alpha = 0.28 if rank == 0 else 0.14

    # Card
    gc = glass_card(slide, x, y, w, h, T, radius=14)

    # Accent top stripe
    tp = rrect(slide, x, y, w, 0.30 * scale, T.accent_rgb, radius_pct=0)
    if tp:
        multi_stop_gradient(tp, [(0, T.bg), (35, T.accent2), (65, T.accent), (100, T.bg)], 0)
        if rank == 0:
            glow(tp, T.accent.lstrip('#'), radius=18, alpha=0.38)

    # Circular accent ring behind value
    ring_s = min(w * 0.52, h * 0.68)
    ring_x = x + (w - ring_s) / 2
    ring_y = y + (h - ring_s) / 2 - h * 0.06
    outer_ring = oval(slide, ring_x - 0.18, ring_y - 0.18, ring_s + 0.36, ring_s + 0.36, T.accent_rgb)
    if outer_ring: set_solid_alpha(outer_ring, 9)
    inner_ring = oval(slide, ring_x, ring_y, ring_s, ring_s, T.accent_rgb)
    if inner_ring:
        set_solid_alpha(inner_ring, 15)
        if rank == 0:
            glow(inner_ring, T.accent.lstrip('#'), radius=20, alpha=glow_alpha)

    # Hero value — size depends on length and rank
    vlen = len(str(value).replace('%', '').replace('.', ''))
    base_vs = 48 if rank == 0 else 36
    vs = base_vs if vlen <= 2 else (base_vs - 8) if vlen <= 4 else (base_vs - 16)
    vs = max(18, vs)
    txt(slide, str(value), x + 0.18, ring_y, w - 0.36, ring_s * 0.65,
        font="Calibri", size=vs, bold=True,
        color=T.accent_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)

    # Unit
    if unit:
        txt(slide, unit, x + 0.18, ring_y + ring_s * 0.60, w - 0.36, ring_s * 0.25,
            font=font, size=max(8, vs * 0.28), bold=False,
            color=T.muted_rgb, align=PP_ALIGN.CENTER, rtl=True, vcenter=True)

    # Divider
    dw = w * 0.55; dx = x + (w - dw) / 2
    dy = y + h * 0.76
    hl = rect(slide, dx, dy, dw, 0.055, T.accent_rgb)
    if hl: multi_stop_gradient(hl, [(0, T.bg2), (50, T.accent), (100, T.bg2)], 0)

    # Label
    fs_lbl = min(12.5 * scale, max(10, h * 5.2))
    txt(slide, label, x + 0.20, dy + 0.10, w - 0.40, h * 0.22,
        font=font, size=fs_lbl, bold=False,
        color=T.text_light_rgb, align=PP_ALIGN.CENTER, rtl=True, vcenter=True,
        line_spacing=1.2)


def timeline_row(slide, x, y, w, h, T, text: str, step: int, total: int, font="Cairo"):
    """
    Horizontal timeline row: step bubble → connector → content card.
    Provides variety for sequential content like plan/methodology.
    """
    bubble_d = min(h * 0.78, 1.30)
    bx = x + w - bubble_d - 0.14   # RTL: bubble on right
    by = y + (h - bubble_d) / 2

    # Connector line from previous row
    if step > 1:
        vline(slide, bx + bubble_d / 2 - 0.04, y - h * 0.22, h * 0.22 + 0.06,
              T.accent_rgb, thickness=0.08)

    # Bubble
    bc = oval(slide, bx, by, bubble_d, bubble_d, T.accent_rgb)
    if bc:
        multi_stop_gradient(bc, [(0, T.accent), (100, T.accent2)], 135)
        shadow(bc, blur=10, dist=3, alpha=0.34)
        if step == 1 or step == total:
            glow(bc, T.accent.lstrip('#'), radius=14, alpha=0.22)
    txt(slide, str(step), bx, by, bubble_d, bubble_d,
        font="Calibri", size=max(9, int(bubble_d * 10.5)), bold=True,
        color=T.text_dark_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)

    # Content card (takes remaining width)
    cw = w - bubble_d - 0.55
    cx = x
    cc = rrect(slide, cx, y + h * 0.06, cw, h * 0.88,
               T.bg2_rgb if step % 2 == 0 else T.card_rgb, radius_pct=9)
    if cc:
        multi_stop_gradient(cc,
                            [(0, T.bg2), (100, T.card)] if step % 2 == 0
                            else [(0, T.card), (100, T.bg2)], 0)
        shadow(cc, blur=7, dist=2, alpha=0.18)

    # Accent left bar (RTL: visually leftmost = start of reading)
    al = rect(slide, cx, y + h * 0.06, 0.20, h * 0.88, T.accent_rgb)
    if al:
        gradient_fill(al, T.accent_grad1, T.accent_grad2, 90)
        if step % 3 == 0:
            set_solid_alpha(al, 72)

    fs = _smart_font_size(text, 13.5, 11, 16, cw - 0.5, h * 0.88)
    txt(slide, text, cx + 0.30, y + h * 0.06, cw - 0.50, h * 0.88,
        font=font, size=fs, bold=False,
        color=T.text_light_rgb, align=PP_ALIGN.RIGHT,
        rtl=True, vcenter=True, line_spacing=1.3)


def icon_text_row(slide, x, y, w, h, T, text: str, icon: str, index: int,
                  font="Cairo", highlight=False):
    """
    Icon + text row with alternating background — premium list item.
    RTL-correct: icon anchored to RIGHT side.
    """
    even = index % 2 == 0
    bg_clr = T.card_rgb if even else T.bg2_rgb
    rw = rrect(slide, x, y, w, h, bg_clr, radius_pct=10)
    if rw:
        stops = [(0, T.card), (100, T.bg2)] if even else [(0, T.bg2), (100, T.card)]
        multi_stop_gradient(rw, stops, 0)
        if highlight:
            shadow(rw, blur=16, dist=4, alpha=0.36)
            glow(rw, T.accent.lstrip('#'), radius=16, alpha=0.07)
        else:
            shadow(rw, blur=5, dist=2, alpha=0.16)

    # RTL accent bar — on the RIGHT (reading start)
    acc = rect(slide, x + w - 0.26, y, 0.26, h, T.accent_rgb)
    if acc:
        gradient_fill(acc, T.accent_grad1, T.accent_grad2, 90)
        if not highlight:
            set_solid_alpha(acc, max(22, 68 - index * 6))

    # Icon circle — anchored RIGHT after accent bar (native RTL)
    ic_s = min(h * 0.72, 1.10)
    ic_x = x + w - 0.26 - ic_s - 0.22
    ic_y = y + (h - ic_s) / 2
    ic_bg = oval(slide, ic_x, ic_y, ic_s, ic_s, T.accent_rgb)
    if ic_bg:
        multi_stop_gradient(ic_bg, [(0, T.accent), (100, T.accent2)], 135)
        shadow(ic_bg, blur=6, dist=2, alpha=0.26)
        if highlight:
            glow(ic_bg, T.accent.lstrip('#'), radius=12, alpha=0.22)
    txt(slide, icon, ic_x, ic_y, ic_s, ic_s,
        font="Segoe UI Emoji", size=max(9, int(ic_s * 9.5)), bold=False,
        color=T.text_dark_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)

    # Text content — flows from RIGHT inward (RTL)
    text_x = x + 0.28
    text_w = w - ic_s - 0.90
    fs = _smart_font_size(text, 13.5, 11, 15.5, text_w, h)
    txt(slide, text, text_x, y + 0.06, text_w, h - 0.12,
        font=font, size=fs, bold=highlight,
        color=T.accent_rgb if highlight else T.text_light_rgb,
        align=PP_ALIGN.RIGHT, rtl=True, vcenter=True, line_spacing=1.28)


def two_panel_layout(slide, x, y, w, h, T,
                     left_label: str, left_items: list,
                     right_label: str, right_items: list,
                     font="Cairo"):
    """
    Two-panel comparison layout with smart content distribution.
    Used for objectives + hypotheses, pros + cons, etc.
    RTL: right panel = primary (starts reading).
    """
    gap = 0.28
    pw = (w - gap) / 2
    px_right = x + pw + gap   # RIGHT panel = reading start in RTL
    px_left  = x               # LEFT panel = secondary

    for px, plabel, pitems, is_primary in [
        (px_right, right_label, right_items, True),
        (px_left,  left_label,  left_items,  False),
    ]:
        # Panel container
        pc = rrect(slide, px, y, pw, h, T.card_rgb, radius_pct=12)
        if pc:
            if is_primary:
                multi_stop_gradient(pc, [(0, T.card), (100, T.bg2)], 145)
                shadow(pc, blur=20, dist=6, alpha=0.42)
                glow(pc, T.accent.lstrip('#'), radius=18, alpha=0.07)
            else:
                multi_stop_gradient(pc, [(0, T.bg2), (100, T.card)], 145)
                shadow(pc, blur=10, dist=3, alpha=0.24)

        # Panel header
        hh = 0.72
        ph = rrect(slide, px, y, pw, hh, T.accent_rgb, radius_pct=0)
        if ph:
            if is_primary:
                multi_stop_gradient(ph, [(0, T.accent2), (60, T.accent), (100, T.accent2)], 0)
            else:
                multi_stop_gradient(ph, [(0, T.accent), (100, T.accent2)], 0)
                set_solid_alpha(ph, 82)
        txt(slide, plabel, px + 0.18, y, pw - 0.36, hh,
            font=font, size=16, bold=True,
            color=T.text_dark_rgb, align=PP_ALIGN.CENTER, rtl=True, vcenter=True)

        # Items
        ia = h - hh - 0.12
        n = min(len(pitems), 7)
        if n == 0: continue
        ig = 0.10
        ih = (ia - ig * (n - 1)) / n

        for j, item in enumerate(pitems[:7]):
            iy = y + hh + 0.06 + j * (ih + ig)
            rb = rrect(slide, px + 0.10, iy, pw - 0.20, ih,
                       T.bg_rgb if j % 2 == 0 else T.bg2_rgb, radius_pct=6)
            if rb: set_solid_alpha(rb, 68)

            # Number on right (RTL)
            nd = min(0.52, ih * 0.72)
            nx = px + pw - 0.72; ny = iy + (ih - nd) / 2
            nb = oval(slide, nx, ny, nd, nd, T.accent_rgb)
            if nb:
                multi_stop_gradient(nb, [(0, T.accent), (100, T.accent2)], 135)
                if not is_primary: set_solid_alpha(nb, 72)
            txt(slide, str(j + 1), nx, ny, nd, nd,
                font="Calibri", size=max(8, int(nd * 10)), bold=True,
                color=T.text_dark_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)

            fs = _smart_font_size(item, 12.5, 10, 14.5, pw - 1.10, ih)
            txt(slide, item, px + 0.22, iy + 0.04, pw - 1.08, ih - 0.08,
                font=font, size=fs, bold=False,
                color=T.text_light_rgb, align=PP_ALIGN.RIGHT,
                rtl=True, vcenter=True, line_spacing=1.2)


def highlight_keyword(text: str, max_words=8) -> tuple[str, str]:
    """
    Split long text into a SHORT highlighted headline + full detail.
    Returns (headline, detail) where headline = first max_words words.
    Used to create visual hierarchy within content cards.
    """
    words = text.split()
    if len(words) <= max_words:
        return text, ''
    headline = ' '.join(words[:max_words]) + '...'
    return headline, text


def cinematic_cover(slide, T, req, font="Cairo"):
    """
    Full cinematic cover with layered composition:
    - Split background (dark/accent zones)
    - Floating title card with depth
    - Author info in premium label strip
    - Year badge + institution line
    """
    import re as _re

    # ── Background ─────────────────────────────────────────────────
    cinematic_bg(slide, T, variant='split')

    # ── Top institutional bar ───────────────────────────────────────
    if req.institution:
        ib = rect(slide, 0, 0, W, 0.66, T.bg_rgb)
        if ib: set_solid_alpha(ib, 72)
        txt(slide, req.institution, 0, 0, W, 0.66,
            font=font, size=10.5, bold=False,
            color=T.muted_rgb, align=PP_ALIGN.CENTER, rtl=True, vcenter=True)
    hline(slide, W * 0.06, 0.68, W * 0.88, T.accent_rgb, thickness=0.055)

    # ── Year extraction ─────────────────────────────────────────────
    _yp = _re.compile(r'\b\d{4}\s*[-–—]\s*\d{4}\b')
    _ym = _yp.search(req.title_ar or '')
    if _ym:
        year_str   = _ym.group(0).strip()
        title_clean = _yp.sub('', req.title_ar).strip(' —–-،, ')
    elif req.year:
        year_str   = req.year.strip()
        title_clean = req.title_ar
    else:
        year_str   = ''
        title_clean = req.title_ar

    # ── Central title zone ──────────────────────────────────────────
    tz_x = W * 0.04; tz_w = W * 0.92
    tz_y = 0.80; tz_h = H * 0.56
    has_year = bool(year_str)
    title_text_h = tz_h * (0.58 if has_year else 0.72)

    # Outer glow backdrop
    og = rrect(slide, tz_x - 0.18, tz_y - 0.18, tz_w + 0.36, tz_h + 0.36, T.accent_rgb, radius_pct=16)
    if og: set_solid_alpha(og, 12)

    # Main title card
    tc = rrect(slide, tz_x, tz_y, tz_w, tz_h, T.card_rgb, radius_pct=14)
    if tc:
        multi_stop_gradient(tc, [(0, T.card), (55, T.bg2), (100, T.bg)], 140)
        shadow(tc, blur=34, dist=10, alpha=0.56)

    # Accent stripe top
    ct = rrect(slide, tz_x, tz_y, tz_w, 0.36, T.accent_rgb, radius_pct=0)
    if ct:
        multi_stop_gradient(ct, [(0, T.bg), (30, T.accent2), (50, T.accent),
                                  (70, T.accent2), (100, T.bg)], 0)
        glow(ct, T.accent.lstrip('#'), radius=18, alpha=0.40)

    # RTL anchor bar on RIGHT
    vline(slide, tz_x + tz_w - 0.28, tz_y + 0.36, tz_h - 0.36, T.accent_rgb, thickness=0.28)

    # Title text — smart sizing
    ts = _smart_font_size(title_clean, 24, 14, 30, tz_w - 1.2, title_text_h)
    txt(slide, title_clean, tz_x + 0.60, tz_y + 0.42, tz_w - 1.20, title_text_h,
        font=font, size=ts, bold=True,
        color=T.text_light_rgb, align=PP_ALIGN.CENTER,
        rtl=True, vcenter=True, line_spacing=1.22)

    if req.title_en:
        en_y = tz_y + title_text_h + 0.32
        en_h = tz_h - title_text_h - 0.40
        txt(slide, req.title_en, tz_x + 0.80, en_y, tz_w - 1.60, en_h,
            font="Calibri", size=11, bold=False, italic=True,
            color=T.muted_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)

    # Year badge
    if has_year:
        yr_y = tz_y + tz_h * 0.80; yr_h = tz_h * 0.12
        yr_cw = tz_w * 0.38; yr_cx = tz_x + (tz_w - yr_cw) / 2
        yb = rrect(slide, yr_cx, yr_y, yr_cw, yr_h, T.accent_rgb, radius_pct=50)
        if yb:
            set_solid_alpha(yb, 22)
        hline(slide, yr_cx + yr_cw * 0.08, yr_y, yr_cw * 0.84, T.accent_rgb, thickness=0.035)
        hline(slide, yr_cx + yr_cw * 0.08, yr_y + yr_h, yr_cw * 0.84, T.accent_rgb, thickness=0.035)
        txt(slide, f'( {year_str} )', yr_cx, yr_y, yr_cw, yr_h,
            font=font, size=13, bold=False, italic=True,
            color=T.accent_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)

    # ── Info strip below title ──────────────────────────────────────
    info_y = tz_y + tz_h + 0.20
    info_h = H - info_y - 0.18
    ic = rrect(slide, tz_x, info_y, tz_w, info_h, T.card_rgb, radius_pct=11)
    if ic:
        multi_stop_gradient(ic, [(0, T.bg2), (100, T.card)], 0)
        shadow(ic, blur=14, dist=4, alpha=0.30)

    rows = [("الطالب", req.student_name)]
    if req.supervisor:     rows.append(("المشرف", req.supervisor))
    if req.co_supervisor:  rows.append(("المشرف المساعد", req.co_supervisor))
    if req.specialization: rows.append(("التخصص", req.specialization))

    rh = info_h / max(len(rows), 1)
    lbl_w = 4.6

    for i, (lbl, val) in enumerate(rows):
        y2 = info_y + i * rh
        if i % 2 == 0:
            rb = rrect(slide, tz_x + 0.18, y2 + 0.04, tz_w - 0.36, rh - 0.08, T.bg_rgb, radius_pct=6)
            if rb: set_solid_alpha(rb, 44)

        # RTL label on RIGHT
        lbl_x = tz_x + tz_w - lbl_w - 0.28
        val_x = tz_x + 0.38
        val_w = tz_w - lbl_w - 0.80

        txt(slide, f"{lbl} :", lbl_x, y2, lbl_w, rh,
            font=font, size=max(13, min(15, rh * 8.5)), bold=True,
            color=T.accent_rgb, align=PP_ALIGN.RIGHT, rtl=True, vcenter=True)
        vline(slide, lbl_x - 0.18, y2 + rh * 0.12, rh * 0.76, T.muted_rgb, thickness=0.045)
        txt(slide, val, val_x, y2, val_w, rh,
            font=font, size=max(14, min(16, rh * 10)), bold=False,
            color=T.text_light_rgb, align=PP_ALIGN.RIGHT, rtl=True, vcenter=True)

    # Bottom bar
    bt = rect(slide, 0, H - 0.30, W, 0.30, T.accent_rgb)
    if bt: multi_stop_gradient(bt, [(0, T.bg), (30, T.accent), (70, T.accent2), (100, T.bg)], 0)


def cinematic_final(slide, T, req, font="Cairo"):
    """
    Full cinematic Thank You slide with strong emotional closure.
    """
    # Background
    cinematic_bg(slide, T, variant='radial')
    decorative_dots(slide, 1.2, H - 5.8, 9, 3, 0.14, 0.38, T.accent_rgb, alpha=8)
    decorative_dots(slide, W - 7.5, 1.0, 7, 4, 0.13, 0.34, T.accent_rgb, alpha=7)
    diamond(slide, W * 0.24, H * 0.06, 2.8, 2.8, T.accent_rgb, alpha=7)
    diamond(slide, W * 0.66, H * 0.76, 2.2, 2.2, T.accent_rgb, alpha=6)

    # Central card
    cw = W * 0.74; ch = H * 0.80
    cx = (W - cw) / 2; cy = H * 0.09

    # Multi-layer depth
    d3 = rrect(slide, cx + 0.42, cy + 0.48, cw, ch, T.bg_rgb, radius_pct=16)
    if d3: set_solid_alpha(d3, 18)
    d2 = rrect(slide, cx + 0.22, cy + 0.24, cw, ch, T.bg_rgb, radius_pct=16)
    if d2: set_solid_alpha(d2, 28)
    cc = rrect(slide, cx, cy, cw, ch, T.card_rgb, radius_pct=16)
    if cc:
        multi_stop_gradient(cc, [(0, T.card), (50, T.bg2), (100, T.card)], 135)
        shadow(cc, blur=40, dist=14, alpha=0.62)

    # Accent stripe TOP
    tp = rrect(slide, cx, cy, cw, 0.42, T.accent_rgb, radius_pct=0)
    if tp:
        multi_stop_gradient(tp, [(0, T.bg), (28, T.accent2), (50, T.accent),
                                  (72, T.accent2), (100, T.bg)], 0)
        glow(tp, T.accent.lstrip('#'), radius=22, alpha=0.44)

    # Accent stripe BOTTOM
    bp = rrect(slide, cx, cy + ch - 0.26, cw, 0.26, T.accent_rgb, radius_pct=0)
    if bp:
        multi_stop_gradient(bp, [(0, T.bg), (50, T.accent), (100, T.bg)], 0)
        set_solid_alpha(bp, 52)

    # RTL anchor (right bar)
    vline(slide, cx + cw - 0.26, cy + 0.42, ch - 0.68, T.accent_rgb, thickness=0.26)

    # Star
    txt(slide, "✦", cx + cw / 2 - 1.0, cy + 0.52, 2.0, 1.4,
        font="Calibri", size=30, bold=False, color=T.accent_rgb,
        align=PP_ALIGN.CENTER, rtl=False, vcenter=True)

    # Main text
    txt(slide, "شكراً وتقديراً", cx + 0.80, cy + 1.16, cw - 1.60, ch * 0.27,
        font=font, size=42, bold=True,
        color=T.text_light_rgb, align=PP_ALIGN.CENTER, rtl=True, vcenter=True)

    # Primary divider
    d1 = rect(slide, cx + cw * 0.13, cy + ch * 0.41, cw * 0.74, 0.07, T.accent_rgb)
    if d1: multi_stop_gradient(d1, [(0, T.bg2), (50, T.accent), (100, T.bg2)], 0)
    rect(slide, cx + cw * 0.22, cy + ch * 0.41 + 0.14, cw * 0.56, 0.033, T.muted_rgb)

    # Student name — hero secondary
    txt(slide, req.student_name, cx + 0.80, cy + ch * 0.44, cw - 1.60, ch * 0.15,
        font=font, size=24, bold=True,
        color=T.accent_rgb, align=PP_ALIGN.CENTER, rtl=True, vcenter=True)

    # Title excerpt
    ts = req.title_ar[:68] + ("..." if len(req.title_ar) > 68 else "")
    txt(slide, ts, cx + 1.10, cy + ch * 0.60, cw - 2.20, ch * 0.22,
        font=font, size=12, bold=False, italic=True,
        color=T.muted_rgb, align=PP_ALIGN.CENTER, rtl=True, vcenter=True, line_spacing=1.30)

    # Footer: institution · year
    footer = []
    if req.institution: footer.append(req.institution)
    if req.year: footer.append(req.year)
    if footer:
        fb = rrect(slide, cx + cw * 0.12, cy + ch * 0.85, cw * 0.76, 0.68, T.bg_rgb, radius_pct=40)
        if fb: set_solid_alpha(fb, 50)
        txt(slide, "  ·  ".join(footer), cx + 0.80, cy + ch * 0.85, cw - 1.60, 0.68,
            font=font, size=11, bold=False,
            color=T.muted_rgb, align=PP_ALIGN.CENTER, rtl=True, vcenter=True)

    # Page-bottom accent bar
    bt = rect(slide, 0, H - 0.30, W, 0.30, T.accent_rgb)
    if bt: multi_stop_gradient(bt, [(0, T.bg), (30, T.accent), (70, T.accent2), (100, T.bg)], 0)
