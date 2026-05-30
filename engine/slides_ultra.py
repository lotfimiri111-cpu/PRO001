"""
Ultra Engine — مذكرتي Pro v20 Ultra
محرك تصميم جديد كلياً بمستوى Canva / Gamma / Beautiful.ai

المبادئ:
- RTL-native rendering
- Visual hierarchy حقيقي
- Layout intelligence — لا تكرار
- Cinematic cover + strong endings
- Hero statistics
- Arabic-first typography
- Smart whitespace
"""
from __future__ import annotations
import math
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Cm, Pt
from lxml import etree
from pptx.oxml.ns import qn

from engine.primitives import (
    W, H, rect, rrect, oval, bg, hline, vline,
    gradient_fill, gradient_rect, shadow, set_solid_alpha,
    multi_stop_gradient, glow, diamond, hexagon, decorative_dots,
    txt, blank_slide, cm, pt,
)
from core.themes import Theme
from core.models import PresentationRequest

_FONT = "Cairo"
def set_font(n): global _FONT; _FONT = n

# ══════════════════════════════════════════════════════════════════════
# SMART TEXT TRUNCATOR
# ══════════════════════════════════════════════════════════════════════
def _trunc(text: str, max_words: int = 18) -> str:
    """تقليص النص لعرض presentation مناسب"""
    if not text:
        return ""
    words = text.split()
    if len(words) <= max_words:
        return text
    return " ".join(words[:max_words]) + "..."

def _font_size(text: str, base: int, max_chars: int = 40) -> int:
    """حجم خط ذكي حسب طول النص"""
    if len(text) <= max_chars:
        return base
    elif len(text) <= max_chars * 1.5:
        return max(int(base * 0.85), 10)
    else:
        return max(int(base * 0.72), 9)

# ══════════════════════════════════════════════════════════════════════
# CINEMATIC BACKGROUND SYSTEM
# ══════════════════════════════════════════════════════════════════════
def _bg_ultra(slide, T, variant='default'):
    """خلفية سينمائية متعددة الطبقات"""
    # Base gradient
    bg(slide, T.bg_rgb)
    gradient_rect(slide, 0, 0, W, H, T.grad1, T.grad2, angle=145)

    if variant == 'cover':
        # طبقات هندسية للغلاف
        gradient_rect(slide, -2, -2, W * 0.55, H + 4, T.bg2, T.bg, angle=160)
        # شريط مائل أيمن
        s1 = rect(slide, W * 0.62, 0, W * 0.45, H, T.accent_rgb)
        if s1:
            multi_stop_gradient(s1, [(0, T.accent), (60, T.bg2), (100, T.bg)], angle=180)
            set_solid_alpha(s1, 18)
        # دوائر خلفية ضخمة
        oval(slide, W * 0.5, -H * 0.3, H * 1.4, H * 1.4, T.bg2_rgb, alpha=55)
        oval(slide, -W * 0.1, H * 0.2, H * 0.8, H * 0.8, T.accent_rgb, alpha=6)
        # نقاط زخرفية
        decorative_dots(slide, 0.6, 1.2, 4, 6, 0.14, 0.38, T.accent_rgb, alpha=14)
        decorative_dots(slide, W - 5.5, H - 4.5, 5, 4, 0.12, 0.32, T.accent_rgb, alpha=9)

    elif variant == 'section':
        # تقسيم عمودي للشرائح الرئيسية
        s = rect(slide, 0, 0, W * 0.38, H, T.bg2_rgb)
        if s:
            gradient_fill(s, T.bg2, T.bg, angle=180)
        oval(slide, W * 0.22, -3, 10, 10, T.accent_rgb, alpha=6)
        decorative_dots(slide, 1.0, H - 3.5, 4, 3, 0.15, 0.35, T.accent_rgb, alpha=11)

    elif variant == 'split':
        # تقسيم أفقي
        r = rect(slide, 0, H * 0.56, W, H * 0.44, T.bg2_rgb)
        if r:
            gradient_fill(r, T.bg2, T.bg, angle=90)
        oval(slide, W - 8, -2, 12, 12, T.accent_rgb, alpha=5)

    elif variant == 'diagonal':
        # خط قطري
        s = rect(slide, W * 0.7, 0, W * 0.35, H, T.bg2_rgb)
        if s:
            gradient_fill(s, T.accent, T.bg2, angle=90)
            set_solid_alpha(s, 12)
        oval(slide, -3, H - 8, 11, 11, T.accent_rgb, alpha=5)
        decorative_dots(slide, W - 7, 1.5, 5, 3, 0.16, 0.4, T.accent_rgb, alpha=8)

    else:  # default
        oval(slide, -3, -3, 10, 10, T.accent_rgb, alpha=5)
        oval(slide, W - 8, H - 7, 12, 12, T.bg2_rgb, alpha=45)
        decorative_dots(slide, 1.2, H - 3.8, 5, 3, 0.15, 0.38, T.accent_rgb, alpha=10)

# ══════════════════════════════════════════════════════════════════════
# RTL-NATIVE HEADER — تصميم أصلي للعربية
# ══════════════════════════════════════════════════════════════════════
HEADER_H = 2.65

def _hdr_ultra(slide, T, title: str, sub: str = '', num: int = None, total: int = 13):
    """هيدر RTL أصلي — العنوان يبدأ من اليمين بشكل طبيعي"""
    # الخلفية
    gradient_rect(slide, 0, 0, W, HEADER_H, T.grad2, T.grad1, angle=140)

    # شريط accent أسفل الهيدر
    al = rect(slide, 0, HEADER_H - 0.18, W, 0.18, T.accent_rgb)
    if al:
        multi_stop_gradient(al, [(0, T.bg), (35, T.accent), (65, T.accent2), (100, T.bg)], 0)

    # شريط accent عمودي — في RTL يكون على اليمين
    av = rect(slide, W - 0.48, 0, 0.48, HEADER_H - 0.18, T.accent_rgb)
    if av:
        gradient_fill(av, T.accent_grad1, T.accent_grad2, 90)

    # رقم الشريحة — دائرة في أقصى اليسار (عكس RTL)
    num_end = 0.5
    if num is not None:
        nb_s = 0.68
        nb_x = 0.8
        nb_y = (HEADER_H - 0.18 - nb_s) / 2
        nb = oval(slide, nb_x, nb_y, nb_s, nb_s, T.accent_rgb)
        if nb:
            gradient_fill(nb, T.accent_grad1, T.accent_grad2, 135)
            shadow(nb, blur=8, dist=2, alpha=0.4)
        txt(slide, str(num), nb_x, nb_y, nb_s, nb_s,
            font="Calibri", size=13, bold=True,
            color=T.text_dark_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)
        txt(slide, f"/{total}", nb_x + nb_s, nb_y + nb_s * 0.32, 0.75, nb_s * 0.4,
            font="Calibri", size=8, bold=False,
            color=T.muted_rgb, align=PP_ALIGN.LEFT, rtl=False, vcenter=True)
        num_end = nb_x + nb_s + 0.85

    # العنوان — من اليمين (RTL) إلى نقطة انتهاء رقم الشريحة
    title_w = W - num_end - 0.7 - 0.55  # 0.7 للشريط الأيمن + padding
    title_fs = _font_size(title, 24, 35)
    txt(slide, title, num_end, 0.18, title_w, HEADER_H * 0.62,
        font=_FONT, size=title_fs, bold=True,
        color=T.text_light_rgb, align=PP_ALIGN.RIGHT,
        rtl=True, vcenter=True, line_spacing=1.1)

    # العنوان الفرعي
    if sub:
        txt(slide, sub, num_end, HEADER_H * 0.6, title_w, HEADER_H * 0.35,
            font=_FONT, size=12, bold=False, italic=True,
            color=T.muted_rgb, align=PP_ALIGN.RIGHT,
            rtl=True, vcenter=True, line_spacing=1.0)

# منطقة المحتوى
CY0 = HEADER_H + 0.32
def _ch(): return H - CY0 - 0.3

# ══════════════════════════════════════════════════════════════════════
# BULLET SYSTEM — RTL native
# ══════════════════════════════════════════════════════════════════════
def _rtl_bullet_list(slide, T, items: list[str], x, y, w, max_items=6,
                     icon='◆', font_size=13, line_h=None):
    """قائمة bullets أصيلة RTL — الأيقونة على اليمين، النص يتدفق لليسار"""
    if not items:
        return
    items = items[:max_items]
    n = len(items)
    if line_h is None:
        available_h = H - y - 0.3
        line_h = min(available_h / n, 1.55)

    for i, item in enumerate(items):
        iy = y + i * line_h
        item_text = _trunc(item, 22)

        # خلفية البطاقة (أكثر شفافية للعناصر الأخيرة)
        alpha = 35 if i % 2 == 0 else 22
        card = rrect(slide, x, iy + 0.06, w, line_h - 0.12, T.card_rgb, radius_pct=7)
        if card:
            set_solid_alpha(card, alpha)

        # شريط accent على اليمين (RTL)
        accent_bar = rect(slide, x + w - 0.18, iy + 0.14, 0.18, line_h - 0.28, T.accent_rgb)
        if accent_bar:
            gradient_fill(accent_bar, T.accent_grad1, T.accent_grad2, 90)

        # رقم أو أيقونة
        num_s = 0.42
        num_x = x + w - num_s - 0.28
        num_circle = oval(slide, num_x, iy + (line_h - num_s) / 2,
                          num_s, num_s, T.accent_rgb)
        if num_circle:
            gradient_fill(num_circle, T.accent_grad1, T.accent_grad2, 135)
        txt(slide, str(i + 1), num_x, iy + (line_h - num_s) / 2, num_s, num_s,
            font="Calibri", size=9, bold=True,
            color=T.text_dark_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)

        # النص — يبدأ من اليمين ويتدفق لليسار
        txt_x = x + 0.28
        txt_w = w - num_s - 0.28 - 0.55
        fs = _font_size(item_text, font_size, 50)
        txt(slide, item_text, txt_x, iy + 0.06, txt_w, line_h - 0.12,
            font=_FONT, size=fs, bold=False,
            color=T.text_light_rgb, align=PP_ALIGN.RIGHT,
            rtl=True, vcenter=True, line_spacing=1.1)

# ══════════════════════════════════════════════════════════════════════
# COVER — سينمائي
# ══════════════════════════════════════════════════════════════════════
def make_cover(prs, req: PresentationRequest, T: Theme):
    slide = blank_slide(prs)
    _bg_ultra(slide, T, 'cover')

    # شريط علوي — المؤسسة
    top = rect(slide, 0, 0, W, 0.52, T.bg2_rgb)
    if top:
        gradient_fill(top, T.bg2, T.bg, angle=0)
        set_solid_alpha(top, 80)
    if req.institution:
        txt(slide, req.institution, 0.8, 0, W - 1.6, 0.52,
            font=_FONT, size=10.5, bold=False,
            color=T.muted_rgb, align=PP_ALIGN.CENTER, rtl=True, vcenter=True)

    # الخط الفاصل الذهبي
    hl_top = rect(slide, 0, 0.52, W, 0.055, T.accent_rgb)
    if hl_top:
        multi_stop_gradient(hl_top, [(0, T.bg), (30, T.accent), (70, T.accent2), (100, T.bg)], 0)

    # ── البطاقة الرئيسية — العنوان ──
    title_len = len(req.title_ar)
    title_h = 6.8 if title_len < 50 else 7.5 if title_len < 80 else 8.2

    card_x, card_w = 1.5, W - 3.0
    card_y = 0.9

    # ظل عميق
    shadow_c = rrect(slide, card_x + 0.2, card_y + 0.25, card_w, title_h,
                     T.bg_rgb, radius_pct=12)
    if shadow_c:
        set_solid_alpha(shadow_c, 55)

    # البطاقة نفسها
    mc = rrect(slide, card_x, card_y, card_w, title_h, T.card_rgb, radius_pct=12)
    if mc:
        multi_stop_gradient(mc, [(0, T.card), (60, T.bg2), (100, T.card)], angle=145)
        shadow(mc, blur=28, dist=8, alpha=0.55)

    # شريط accent علوي للبطاقة
    ct = rrect(slide, card_x, card_y, card_w, 0.32, T.accent_rgb, radius_pct=0)
    if ct:
        multi_stop_gradient(ct, [(0, T.accent), (50, T.accent2), (100, T.accent)], 0)

    # شريط accent أيمن (RTL)
    vline(slide, card_x + card_w - 0.2, card_y + 0.32,
          title_h - 0.32, T.accent_rgb, thickness=0.2)

    # أيقونة المذكرة
    icon_y = card_y + 0.55
    icon_circle = oval(slide, card_x + card_w - 2.8, icon_y, 1.5, 1.5, T.accent_rgb)
    if icon_circle:
        gradient_fill(icon_circle, T.accent_grad1, T.accent_grad2, 135)
        shadow(icon_circle, blur=16, dist=4, alpha=0.45)
    txt(slide, "🎓", card_x + card_w - 2.8, icon_y, 1.5, 1.5,
        font="Calibri", size=26, bold=False,
        color=T.text_dark_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)

    # العنوان — hero typography
    title_fs = 26 if title_len < 40 else 21 if title_len < 65 else 17
    txt(slide, req.title_ar,
        card_x + 0.45, card_y + 0.38, card_w - 2.0, title_h * 0.62,
        font=_FONT, size=title_fs, bold=True,
        color=T.text_light_rgb, align=PP_ALIGN.RIGHT,
        rtl=True, vcenter=True, line_spacing=1.25)

    # العنوان الفرعي بالفرنسية
    if req.title_en:
        txt(slide, req.title_en,
            card_x + 0.45, card_y + title_h * 0.62, card_w - 1.0, title_h * 0.2,
            font="Calibri", size=11, bold=False, italic=True,
            color=T.muted_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)

    # خط فاصل داخل البطاقة
    divider = rect(slide, card_x + card_w * 0.08, card_y + title_h * 0.84,
                   card_w * 0.84, 0.045, T.accent_rgb)
    if divider:
        multi_stop_gradient(divider, [(0, T.bg2), (50, T.accent), (100, T.bg2)], 0)

    # ── بطاقة المعلومات ──
    info_y = card_y + title_h + 0.22
    info_h = H - 0.38 - info_y
    if info_h > 0.5:
        rows = []
        if req.student_name: rows.append(("👤", "اسم الطالب", req.student_name))
        if req.supervisor: rows.append(("👨‍🏫", "المشرف", req.supervisor))
        if req.specialization: rows.append(("📚", "التخصص", req.specialization))
        if req.year: rows.append(("📅", "السنة", req.year))

        ic = rrect(slide, card_x, info_y, card_w, info_h, T.card_rgb, radius_pct=10)
        if ic:
            multi_stop_gradient(ic, [(0, T.bg2), (100, T.card)], 135)
            shadow(ic, blur=14, dist=4, alpha=0.35)
        vline(slide, card_x + card_w - 0.16, info_y, info_h, T.accent_rgb, thickness=0.16)

        rh = info_h / max(len(rows), 1)
        for i, (icon, lbl, val) in enumerate(rows):
            ry = info_y + i * rh
            # خط فاصل بين الصفوف
            if i > 0:
                hline(slide, card_x + 0.3, ry, card_w - 0.46, T.muted_rgb, thickness=0.03)

            # الأيقونة
            ic_s = min(rh * 0.55, 0.65)
            ic_y2 = ry + (rh - ic_s) / 2
            ci = oval(slide, card_x + card_w - ic_s - 0.38, ic_y2, ic_s, ic_s, T.accent_rgb)
            if ci:
                gradient_fill(ci, T.accent_grad1, T.accent_grad2, 135)
                set_solid_alpha(ci, 70)
            txt(slide, icon, card_x + card_w - ic_s - 0.38, ic_y2, ic_s, ic_s,
                font="Calibri", size=int(ic_s * 16), bold=False,
                color=T.text_dark_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)

            # التسمية
            lbl_fs = max(9, min(11, int(rh * 7)))
            txt(slide, f"{lbl}:", card_x + card_w - ic_s - 3.6, ry, 2.9, rh,
                font=_FONT, size=lbl_fs, bold=True,
                color=T.accent_rgb, align=PP_ALIGN.RIGHT, rtl=True, vcenter=True)

            # الفاصل العمودي
            vline(slide, card_x + card_w - ic_s - 3.72, ry + rh * 0.15,
                  rh * 0.7, T.muted_rgb, thickness=0.04)

            # القيمة
            val_fs = max(10, min(13, int(rh * 8.5)))
            txt(slide, val, card_x + 0.4, ry, card_w - ic_s - 4.1, rh,
                font=_FONT, size=val_fs, bold=False,
                color=T.text_light_rgb, align=PP_ALIGN.RIGHT, rtl=True, vcenter=True)

    # شريط سفلي
    bot = rect(slide, 0, H - 0.28, W, 0.28, T.accent_rgb)
    if bot:
        multi_stop_gradient(bot, [(0, T.bg), (40, T.accent), (60, T.accent2), (100, T.bg)], 0)

    return slide

# ══════════════════════════════════════════════════════════════════════
# INTRO — تخطيط section مقسّم
# ══════════════════════════════════════════════════════════════════════
def make_intro(prs, req: PresentationRequest, T: Theme):
    slide = blank_slide(prs)
    _bg_ultra(slide, T, 'section')
    _hdr_ultra(slide, T, "مقدمة البحث", "نظرة عامة على الدراسة", num=1,
               total=getattr(req, '_total_slides', 13))
    CY = CY0; CH = _ch()

    items = []
    if req.intro_overview: items.append(("📖", "نظرة عامة", req.intro_overview))
    if req.intro_approach: items.append(("🔬", "المنهج المتبع", req.intro_approach))
    if not items:
        return slide

    n = len(items)
    gap = 0.45
    col_w = (W - 2.2 - gap * (n - 1)) / n
    CARD_H = min(CH * 0.88, 11.0)
    card_y = CY + (CH - CARD_H) / 2

    for i, (icon, lbl, val) in enumerate(items[:2]):
        x = 1.1 + i * (col_w + gap)

        # البطاقة
        cc = rrect(slide, x, card_y, col_w, CARD_H, T.card_rgb, radius_pct=12)
        if cc:
            multi_stop_gradient(cc, [(0, T.card), (100, T.bg2)], 150)
            shadow(cc, blur=22, dist=7, alpha=0.5)

        # شريط accent أعلى البطاقة (RTL = يمين)
        ct = rrect(slide, x + col_w - 0.2, card_y, 0.2, CARD_H, T.accent_rgb, radius_pct=0)
        if ct:
            gradient_fill(ct, T.accent_grad1, T.accent_grad2, 90)

        # دائرة الأيقونة
        ic_s = 1.6
        ic_x = x + (col_w - ic_s) / 2
        ic_y = card_y + 0.45
        ic_bg = oval(slide, ic_x, ic_y, ic_s, ic_s, T.bg2_rgb)
        if ic_bg:
            gradient_fill(ic_bg, T.accent_grad1, T.bg2, 135)
            shadow(ic_bg, blur=14, dist=4, alpha=0.4)
        txt(slide, icon, ic_x, ic_y, ic_s, ic_s,
            font="Calibri", size=28, bold=False,
            color=T.text_light_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)

        # التسمية
        txt(slide, lbl, x + 0.25, card_y + ic_y - card_y + ic_s + 0.2,
            col_w - 0.5, 0.7,
            font=_FONT, size=15, bold=True,
            color=T.accent_rgb, align=PP_ALIGN.CENTER, rtl=True, vcenter=True)

        # الفاصل
        dv = rect(slide, x + col_w * 0.15,
                  card_y + ic_y - card_y + ic_s + 0.98, col_w * 0.7, 0.05, T.accent_rgb)
        if dv:
            gradient_fill(dv, T.bg2, T.accent, 0)

        # النص
        text_y = card_y + ic_s + 1.42
        text_h = CARD_H - (text_y - card_y) - 0.4
        txt_content = _trunc(val, 40)
        txt(slide, txt_content, x + 0.3, text_y, col_w - 0.6, text_h,
            font=_FONT, size=12, bold=False,
            color=T.text_light_rgb, align=PP_ALIGN.RIGHT,
            rtl=True, vcenter=False, line_spacing=1.35)

    return slide

# ══════════════════════════════════════════════════════════════════════
# PLAN (خطة الدراسة) — تخطيط timeline عمودي
# ══════════════════════════════════════════════════════════════════════
def make_plan(prs, req: PresentationRequest, T: Theme):
    slide = blank_slide(prs)
    _bg_ultra(slide, T, 'diagonal')
    _hdr_ultra(slide, T, "خطة الدراسة", "هيكل الأطروحة", num=2,
               total=getattr(req, '_total_slides', 13))
    CY = CY0; CH = _ch()

    chapters = req.chapters[:6]
    if not chapters:
        return slide

    n = len(chapters)

    if n <= 3:
        # تخطيط أفقي للفصول القليلة
        gap = 0.4
        col_w = (W - 2.0 - gap * (n - 1)) / n
        CARD_H = min(CH * 0.82, 10.0)
        cy = CY + (CH - CARD_H) / 2

        for i, ch in enumerate(chapters):
            x = 1.0 + i * (col_w + gap)

            # البطاقة
            c = rrect(slide, x, cy, col_w, CARD_H, T.card_rgb, radius_pct=12)
            if c:
                multi_stop_gradient(c, [(0, T.card), (100, T.bg2)], 145)
                shadow(c, blur=20, dist=6, alpha=0.48)

            # رقم الفصل — hero number
            nb_s = 2.0
            nb_x = x + (col_w - nb_s) / 2
            nb_c = oval(slide, nb_x, cy + 0.5, nb_s, nb_s, T.accent_rgb)
            if nb_c:
                gradient_fill(nb_c, T.accent_grad1, T.accent_grad2, 135)
                shadow(nb_c, blur=16, dist=5, alpha=0.45)
            txt(slide, str(i + 1), nb_x, cy + 0.5, nb_s, nb_s,
                font="Calibri", size=28, bold=True,
                color=T.text_dark_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)

            # العنوان
            title_text = _trunc(ch.title, 15)
            txt(slide, title_text, x + 0.2, cy + nb_s + 1.1, col_w - 0.4, CARD_H - nb_s - 1.5,
                font=_FONT, size=_font_size(title_text, 13, 25), bold=True,
                color=T.text_light_rgb, align=PP_ALIGN.CENTER,
                rtl=True, vcenter=True, line_spacing=1.2)

    else:
        # تخطيط timeline للفصول الكثيرة
        ITEM_H = min((CH - 0.4) / n, 2.4)
        timeline_x = W * 0.5  # خط الوقت في المنتصف
        card_w = W * 0.42

        for i, ch in enumerate(chapters):
            iy = CY + i * ITEM_H + 0.2

            # الفصول الزوجية على اليمين (RTL = البداية)، الفردية على اليسار
            if i % 2 == 0:
                cx = W * 0.52 + 0.3
            else:
                cx = W * 0.08

            c = rrect(slide, cx, iy, card_w, ITEM_H - 0.2, T.card_rgb, radius_pct=8)
            if c:
                multi_stop_gradient(c, [(0, T.card), (100, T.bg2)], 135)
                shadow(c, blur=14, dist=4, alpha=0.42)

            # النقطة على خط الوقت
            dot_s = 0.45
            dot = oval(slide, timeline_x - dot_s / 2, iy + (ITEM_H - dot_s) / 2 - 0.1,
                       dot_s, dot_s, T.accent_rgb)
            if dot:
                gradient_fill(dot, T.accent_grad1, T.accent_grad2, 135)
                shadow(dot, blur=8, dist=2, alpha=0.5)

            # الرقم
            txt(slide, str(i + 1), cx + (card_w - 0.5) * (0 if i % 2 == 0 else 1),
                iy, 0.5, ITEM_H - 0.2,
                font="Calibri", size=11, bold=True,
                color=T.accent_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)

            # العنوان
            txt_x = cx + (0.55 if i % 2 == 0 else 0.1)
            txt(slide, _trunc(ch.title, 12), txt_x, iy, card_w - 0.65, ITEM_H - 0.2,
                font=_FONT, size=12, bold=True,
                color=T.text_light_rgb,
                align=PP_ALIGN.RIGHT if i % 2 == 0 else PP_ALIGN.RIGHT,
                rtl=True, vcenter=True, line_spacing=1.1)

        # خط الوقت العمودي
        vline(slide, timeline_x, CY + 0.2, n * ITEM_H - 0.2, T.accent_rgb, thickness=0.06)

    return slide

# ══════════════════════════════════════════════════════════════════════
# PROBLEM (الإشكالية) — تخطيط hero مع focal point قوي
# ══════════════════════════════════════════════════════════════════════
def make_problem(prs, req: PresentationRequest, T: Theme):
    slide = blank_slide(prs)
    _bg_ultra(slide, T, 'split')
    _hdr_ultra(slide, T, "إشكالية البحث", "التساؤل الجوهري", num=3,
               total=getattr(req, '_total_slides', 13))
    CY = CY0; CH = _ch()

    # السؤال الرئيسي — hero element
    if req.main_question:
        q_h = CH * 0.38
        q_card = rrect(slide, 1.0, CY + 0.1, W - 2.0, q_h, T.accent_rgb, radius_pct=12)
        if q_card:
            multi_stop_gradient(q_card, [(0, T.accent_grad1), (100, T.accent_grad2)], 135)
            shadow(q_card, blur=24, dist=7, alpha=0.5)

        # أيقونة "؟" ضخمة
        txt(slide, "؟", W - 3.5, CY + 0.1, 2.5, q_h,
            font="Calibri", size=52, bold=True,
            color=T.text_dark_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)
        set_solid_alpha(
            slide.shapes[-1], 18)

        txt(slide, req.main_question, 1.3, CY + 0.1, W - 5.0, q_h,
            font=_FONT, size=_font_size(req.main_question, 17, 60), bold=True,
            color=T.text_dark_rgb, align=PP_ALIGN.RIGHT,
            rtl=True, vcenter=True, line_spacing=1.2)

    # الفاصل
    sep_y = CY + CH * 0.42
    hline(slide, 1.2, sep_y, W - 2.4, T.accent_rgb, thickness=0.055)

    # التساؤلات الفرعية — على عمودين
    sub_y = sep_y + 0.25
    sub_h = H - sub_y - 0.35

    sub_qs = req.sub_questions[:4]
    if sub_qs:
        # العنوان
        txt(slide, "التساؤلات الفرعية:", 1.2, sub_y, W - 2.4, 0.55,
            font=_FONT, size=13, bold=True,
            color=T.accent_rgb, align=PP_ALIGN.RIGHT, rtl=True, vcenter=True)

        _rtl_bullet_list(slide, T, sub_qs, 1.2, sub_y + 0.62,
                         W - 2.4, max_items=4,
                         font_size=12,
                         line_h=(sub_h - 0.65) / max(len(sub_qs), 1))
    elif req.main_problem:
        txt(slide, _trunc(req.main_problem, 50), 1.2, sub_y, W - 2.4, sub_h,
            font=_FONT, size=13, bold=False,
            color=T.text_light_rgb, align=PP_ALIGN.RIGHT,
            rtl=True, vcenter=False, line_spacing=1.35)

    return slide

# ══════════════════════════════════════════════════════════════════════
# OBJECTIVES — تخطيط grid مع visual weight
# ══════════════════════════════════════════════════════════════════════
def make_objectives(prs, req: PresentationRequest, T: Theme):
    slide = blank_slide(prs)
    _bg_ultra(slide, T, 'default')
    _hdr_ultra(slide, T, "أهداف وفرضيات البحث", "", num=4,
               total=getattr(req, '_total_slides', 13))
    CY = CY0; CH = _ch()

    has_obj = bool(req.objectives)
    has_hyp = bool(req.hypotheses)

    if has_obj and has_hyp:
        # عمودان متوازيان
        col_w = (W - 3.0) / 2
        _panel(slide, T, "🎯 الأهداف", req.objectives[:5],
               1.0, CY, col_w, CH - 0.1)
        _panel(slide, T, "💡 الفرضيات", req.hypotheses[:4],
               1.0 + col_w + 1.0, CY, col_w, CH - 0.1)
    elif has_obj:
        _rtl_bullet_list(slide, T, req.objectives[:6], 1.2, CY + 0.1,
                         W - 2.4, max_items=6, font_size=13)
    elif has_hyp:
        _rtl_bullet_list(slide, T, req.hypotheses[:5], 1.2, CY + 0.1,
                         W - 2.4, max_items=5, font_size=13)

    return slide


def _panel(slide, T, title, items, x, y, w, h):
    """لوح محتوى مستقل"""
    # بطاقة الخلفية
    c = rrect(slide, x, y, w, h, T.card_rgb, radius_pct=10)
    if c:
        multi_stop_gradient(c, [(0, T.card), (100, T.bg2)], 145)
        shadow(c, blur=18, dist=5, alpha=0.45)

    # شريط العنوان
    hdr = rrect(slide, x, y, w, 0.72, T.accent_rgb, radius_pct=0)
    if hdr:
        gradient_fill(hdr, T.accent_grad1, T.accent_grad2, 0)
    txt(slide, title, x + 0.2, y, w - 0.4, 0.72,
        font=_FONT, size=13, bold=True,
        color=T.text_dark_rgb, align=PP_ALIGN.RIGHT,
        rtl=True, vcenter=True)

    # القائمة
    items_h = h - 0.82
    _rtl_bullet_list(slide, T, items, x + 0.15, y + 0.78,
                     w - 0.3, max_items=6, font_size=12,
                     line_h=items_h / max(len(items), 1))

# ══════════════════════════════════════════════════════════════════════
# IMPORTANCE — تخطيط مغاير (cards horizontal + hero icon)
# ══════════════════════════════════════════════════════════════════════
def make_importance(prs, req: PresentationRequest, T: Theme):
    slide = blank_slide(prs)
    _bg_ultra(slide, T, 'section')
    _hdr_ultra(slide, T, "أهمية البحث", "المساهمة العلمية والتطبيقية", num=5,
               total=getattr(req, '_total_slides', 13))
    CY = CY0; CH = _ch()

    items = req.importance[:5]
    if not items:
        return slide

    # تخطيط: icon hero على اليمين + قائمة على اليسار
    hero_w = W * 0.3
    list_x = 1.0
    list_w = W - hero_w - 2.2

    # hero visual
    hero_x = W - hero_w - 0.8
    hc = rrect(slide, hero_x, CY, hero_w, CH, T.card_rgb, radius_pct=14)
    if hc:
        multi_stop_gradient(hc, [(0, T.accent_grad1), (100, T.bg2)], 145)
        shadow(hc, blur=20, dist=6, alpha=0.45)
        set_solid_alpha(hc, 85)

    txt(slide, "⭐", hero_x, CY + CH * 0.1, hero_w, CH * 0.4,
        font="Calibri", size=52, bold=False,
        color=T.text_light_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)
    txt(slide, "الأهمية", hero_x, CY + CH * 0.5, hero_w, CH * 0.2,
        font=_FONT, size=15, bold=True,
        color=T.accent_rgb, align=PP_ALIGN.CENTER, rtl=True, vcenter=True)
    txt(slide, "العلمية والتطبيقية", hero_x, CY + CH * 0.68, hero_w, CH * 0.2,
        font=_FONT, size=11, bold=False,
        color=T.muted_rgb, align=PP_ALIGN.CENTER, rtl=True, vcenter=True)

    # القائمة
    line_h = CH / max(len(items), 1)
    _rtl_bullet_list(slide, T, items, list_x, CY, list_w,
                     max_items=5, font_size=13, line_h=line_h)

    return slide

# ══════════════════════════════════════════════════════════════════════
# METHODOLOGY — تخطيط horizontal steps
# ══════════════════════════════════════════════════════════════════════
def make_methodology(prs, req: PresentationRequest, T: Theme):
    slide = blank_slide(prs)
    _bg_ultra(slide, T, 'diagonal')
    _hdr_ultra(slide, T, "المنهجية والعينة", "الإطار المنهجي للبحث", num=6,
               total=getattr(req, '_total_slides', 13))
    CY = CY0; CH = _ch()

    # محتوى المنهج — hero text
    if req.methodology:
        meth_h = CH * 0.35
        mc = rrect(slide, 1.0, CY + 0.1, W - 2.0, meth_h, T.card_rgb, radius_pct=10)
        if mc:
            multi_stop_gradient(mc, [(0, T.card), (100, T.bg2)], 135)
            shadow(mc, blur=18, dist=5, alpha=0.45)
        vline(slide, W - 1.25, CY + 0.1, meth_h, T.accent_rgb, thickness=0.22)
        txt(slide, req.methodology, 1.3, CY + 0.1, W - 2.8, meth_h,
            font=_FONT, size=13, bold=False,
            color=T.text_light_rgb, align=PP_ALIGN.RIGHT,
            rtl=True, vcenter=True, line_spacing=1.3)

    # بطاقات العينة
    details = []
    if req.sample_type: details.append(("👥", "نوع العينة", req.sample_type))
    if req.sample_size: details.append(("📊", "حجم العينة", req.sample_size))
    if req.tool:        details.append(("🔧", "أداة الجمع", req.tool))
    if getattr(req, 'spatial_scope', ''): details.append(("📍", "المجال المكاني", req.spatial_scope))

    if details:
        det_y = CY + CH * 0.4 + 0.1
        det_h = H - det_y - 0.35
        n = len(details[:4])
        gap = 0.35
        cw = (W - 2.0 - gap * (n - 1)) / n

        for i, (icon, lbl, val) in enumerate(details[:4]):
            cx = 1.0 + i * (cw + gap)
            c = rrect(slide, cx, det_y, cw, det_h, T.card_rgb, radius_pct=10)
            if c:
                multi_stop_gradient(c, [(0, T.bg2), (100, T.card)], 145)
                shadow(c, blur=14, dist=4, alpha=0.4)
            vline(slide, cx + cw - 0.16, det_y, det_h, T.accent_rgb, thickness=0.16)

            ic_s = 0.7
            ic_x = cx + (cw - ic_s) / 2
            ci = oval(slide, ic_x, det_y + 0.2, ic_s, ic_s, T.accent_rgb)
            if ci:
                gradient_fill(ci, T.accent_grad1, T.accent_grad2, 135)
                set_solid_alpha(ci, 80)
            txt(slide, icon, ic_x, det_y + 0.2, ic_s, ic_s,
                font="Calibri", size=14, bold=False,
                color=T.text_dark_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)

            txt(slide, lbl, cx + 0.15, det_y + ic_s + 0.28, cw - 0.3, 0.55,
                font=_FONT, size=10, bold=True,
                color=T.accent_rgb, align=PP_ALIGN.CENTER, rtl=True, vcenter=True)
            txt(slide, _trunc(val, 12), cx + 0.15, det_y + ic_s + 0.9, cw - 0.3,
                det_h - ic_s - 1.1,
                font=_FONT, size=11, bold=False,
                color=T.text_light_rgb, align=PP_ALIGN.CENTER,
                rtl=True, vcenter=True, line_spacing=1.2)

    return slide

# ══════════════════════════════════════════════════════════════════════
# KPI — Hero Statistics (أقوى شريحة بصرياً)
# ══════════════════════════════════════════════════════════════════════
def make_kpi(prs, req: PresentationRequest, T: Theme):
    slide = blank_slide(prs)
    _bg_ultra(slide, T, 'cover')
    _hdr_ultra(slide, T, "النتائج الإحصائية", "مؤشرات الأداء الرئيسية", num=7,
               total=getattr(req, '_total_slides', 13))
    CY = CY0; CH = _ch()

    stats = req.stats[:6]
    if not stats:
        return slide

    n = len(stats)

    if n <= 3:
        # Hero layout — أرقام ضخمة
        gap = 0.5
        cw = (W - 2.0 - gap * (n - 1)) / n
        CARD_H = CH * 0.86

        for i, stat in enumerate(stats):
            cx = 1.0 + i * (cw + gap)
            cy2 = CY + (CH - CARD_H) / 2

            # ظل
            sh = rrect(slide, cx + 0.2, cy2 + 0.2, cw, CARD_H, T.bg_rgb, radius_pct=14)
            if sh:
                set_solid_alpha(sh, 55)

            # البطاقة
            c = rrect(slide, cx, cy2, cw, CARD_H, T.card_rgb, radius_pct=14)
            if c:
                multi_stop_gradient(c, [(0, T.card), (50, T.bg2), (100, T.card)], 145)
                shadow(c, blur=26, dist=8, alpha=0.55)

            # شريط accent أعلى
            ct = rrect(slide, cx, cy2, cw, 0.28, T.accent_rgb, radius_pct=0)
            if ct:
                gradient_fill(ct, T.accent_grad1, T.accent_grad2, 0)

            # الرقم — hero
            val_text = str(stat.value)
            val_fs = 42 if len(val_text) <= 5 else 32 if len(val_text) <= 8 else 24
            txt(slide, val_text, cx + 0.2, cy2 + 0.4, cw - 0.4, CARD_H * 0.52,
                font="Calibri", size=val_fs, bold=True,
                color=T.accent_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)

            # خط فاصل
            dv = rect(slide, cx + cw * 0.12, cy2 + CARD_H * 0.6,
                      cw * 0.76, 0.05, T.accent_rgb)
            if dv:
                multi_stop_gradient(dv, [(0, T.bg2), (50, T.accent), (100, T.bg2)], 0)

            # التسمية
            lbl_text = _trunc(str(stat.label), 18)
            txt(slide, lbl_text, cx + 0.2, cy2 + CARD_H * 0.64, cw - 0.4,
                CARD_H * 0.28,
                font=_FONT, size=_font_size(lbl_text, 13, 20), bold=False,
                color=T.text_light_rgb, align=PP_ALIGN.CENTER,
                rtl=True, vcenter=True, line_spacing=1.15)

    else:
        # Grid layout للأرقام الكثيرة
        cols = 3
        rows = math.ceil(n / cols)
        gap_x, gap_y = 0.4, 0.4
        cw = (W - 2.0 - gap_x * (cols - 1)) / cols
        CARD_H = (CH - gap_y * (rows - 1)) / rows

        for i, stat in enumerate(stats):
            col = i % cols
            row = i // cols
            cx = 1.0 + col * (cw + gap_x)
            cy2 = CY + row * (CARD_H + gap_y)

            c = rrect(slide, cx, cy2, cw, CARD_H, T.card_rgb, radius_pct=10)
            if c:
                multi_stop_gradient(c, [(0, T.card), (100, T.bg2)], 135)
                shadow(c, blur=16, dist=4, alpha=0.45)
            vline(slide, cx + cw - 0.16, cy2, CARD_H, T.accent_rgb, thickness=0.16)

            val_fs = 26 if len(str(stat.value)) <= 6 else 20
            txt(slide, str(stat.value), cx + 0.2, cy2 + 0.12, cw - 0.5, CARD_H * 0.55,
                font="Calibri", size=val_fs, bold=True,
                color=T.accent_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)

            txt(slide, _trunc(str(stat.label), 16), cx + 0.15, cy2 + CARD_H * 0.6,
                cw - 0.4, CARD_H * 0.32,
                font=_FONT, size=11, bold=False,
                color=T.text_light_rgb, align=PP_ALIGN.CENTER,
                rtl=True, vcenter=True, line_spacing=1.1)

    return slide

# ══════════════════════════════════════════════════════════════════════
# RESULTS — تخطيط asymmetric مع visual flow
# ══════════════════════════════════════════════════════════════════════
def make_results(prs, req: PresentationRequest, T: Theme):
    slide = blank_slide(prs)
    _bg_ultra(slide, T, 'default')
    _hdr_ultra(slide, T, "النتائج الرئيسية", "أبرز ما توصّلت إليه الدراسة", num=8,
               total=getattr(req, '_total_slides', 13))
    CY = CY0; CH = _ch()

    results = req.main_results[:6]
    if not results:
        return slide

    line_h = CH / max(len(results), 1)
    _rtl_bullet_list(slide, T, results, 1.0, CY, W - 2.0,
                     max_items=6, font_size=13, line_h=line_h)

    return slide

# ══════════════════════════════════════════════════════════════════════
# CONCLUSION — تخطيط dramatic مع focus
# ══════════════════════════════════════════════════════════════════════
def make_conclusion(prs, req: PresentationRequest, T: Theme):
    slide = blank_slide(prs)
    _bg_ultra(slide, T, 'cover')
    _hdr_ultra(slide, T, "الخاتمة والاستنتاجات", "", num=9,
               total=getattr(req, '_total_slides', 13))
    CY = CY0; CH = _ch()

    if req.general_conclusion:
        conc_h = CH * 0.48
        cc = rrect(slide, 1.0, CY + 0.1, W - 2.0, conc_h, T.card_rgb, radius_pct=12)
        if cc:
            multi_stop_gradient(cc, [(0, T.accent_grad1), (60, T.bg2), (100, T.card)], 135)
            shadow(cc, blur=24, dist=7, alpha=0.52)

        # اقتباس زخرفي
        txt(slide, '"', W - 2.5, CY + 0.1, 2.0, conc_h,
            font="Calibri", size=72, bold=True,
            color=T.accent_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)
        set_solid_alpha(slide.shapes[-1], 20)

        vline(slide, W - 1.28, CY + 0.1, conc_h, T.accent_rgb, thickness=0.22)
        txt(slide, _trunc(req.general_conclusion, 60),
            1.3, CY + 0.1, W - 2.8, conc_h,
            font=_FONT, size=_font_size(req.general_conclusion, 15, 80), bold=False,
            color=T.text_light_rgb, align=PP_ALIGN.RIGHT,
            rtl=True, vcenter=True, line_spacing=1.35)

    # التوصيات
    if req.recommendations:
        rec_y = CY + CH * 0.52 + 0.15
        txt(slide, "التوصيات:", 1.2, rec_y, W - 2.4, 0.55,
            font=_FONT, size=13, bold=True,
            color=T.accent_rgb, align=PP_ALIGN.RIGHT, rtl=True, vcenter=True)
        rec_h = H - rec_y - 0.9 - 0.35
        _rtl_bullet_list(slide, T, req.recommendations[:4],
                         1.2, rec_y + 0.62, W - 2.4,
                         max_items=4, font_size=12,
                         line_h=rec_h / max(len(req.recommendations[:4]), 1))

    return slide

# ══════════════════════════════════════════════════════════════════════
# RECOMMENDATIONS
# ══════════════════════════════════════════════════════════════════════
def make_recommendations(prs, req: PresentationRequest, T: Theme):
    slide = blank_slide(prs)
    _bg_ultra(slide, T, 'section')
    _hdr_ultra(slide, T, "التوصيات والاقتراحات", "", num=10,
               total=getattr(req, '_total_slides', 13))
    CY = CY0; CH = _ch()

    items = req.recommendations[:5]
    if items:
        line_h = CH / max(len(items), 1)
        _rtl_bullet_list(slide, T, items, 1.0, CY, W - 2.0,
                         max_items=5, font_size=13, line_h=line_h)
    return slide

# ══════════════════════════════════════════════════════════════════════
# FUTURE WORK
# ══════════════════════════════════════════════════════════════════════
def make_future(prs, req: PresentationRequest, T: Theme):
    slide = blank_slide(prs)
    _bg_ultra(slide, T, 'diagonal')
    _hdr_ultra(slide, T, "آفاق البحث المستقبلية", "", num=11,
               total=getattr(req, '_total_slides', 13))
    CY = CY0; CH = _ch()

    items = req.future_work[:5]
    if items:
        line_h = CH / max(len(items), 1)
        _rtl_bullet_list(slide, T, items, 1.0, CY, W - 2.0,
                         max_items=5, font_size=13, line_h=line_h)
    return slide

# ══════════════════════════════════════════════════════════════════════
# REFERENCES
# ══════════════════════════════════════════════════════════════════════
def make_references(prs, req: PresentationRequest, T: Theme):
    slide = blank_slide(prs)
    _bg_ultra(slide, T, 'default')
    _hdr_ultra(slide, T, "المراجع والمصادر", "", num=12,
               total=getattr(req, '_total_slides', 13))
    CY = CY0; CH = _ch()

    refs = req.references[:8]
    if not refs:
        return slide

    line_h = CH / max(len(refs), 1)
    for i, ref in enumerate(refs):
        ry = CY + i * line_h
        # شريط جانبي
        bar = rect(slide, W - 0.22, ry + 0.1, 0.22, line_h - 0.2, T.accent_rgb)
        if bar:
            gradient_fill(bar, T.accent_grad1, T.accent_grad2, 90)
            set_solid_alpha(bar, 60)

        # رقم
        txt(slide, f"[{i+1}]", W - 1.8, ry, 1.4, line_h,
            font="Calibri", size=11, bold=True,
            color=T.accent_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)

        # النص
        fs = _font_size(ref, 11, 80)
        txt(slide, _trunc(ref, 30), 1.0, ry, W - 3.2, line_h,
            font=_FONT, size=fs, bold=False,
            color=T.text_light_rgb, align=PP_ALIGN.RIGHT,
            rtl=True, vcenter=True, line_spacing=1.15)

    return slide

# ══════════════════════════════════════════════════════════════════════
# THANK YOU — خاتمة سينمائية قوية
# ══════════════════════════════════════════════════════════════════════
def make_thankyou(prs, req: PresentationRequest, T: Theme):
    slide = blank_slide(prs)
    _bg_ultra(slide, T, 'cover')

    # طبقات خلفية إضافية للعمق
    oval(slide, W * 0.3, H * 0.1, H * 1.1, H * 1.1, T.accent_rgb, alpha=7)
    oval(slide, -H * 0.2, -H * 0.1, H * 0.9, H * 0.9, T.bg2_rgb, alpha=50)

    # دوائر متحدة المركز — visual focal point
    for r, a in [(9, 4), (7, 6), (5.5, 8), (4, 10)]:
        oval(slide, W / 2 - r / 2, H / 2 - r / 2, r, r, T.accent_rgb, alpha=a)

    # الخط العلوي
    top_bar = rect(slide, 0, 0, W, 0.44, T.bg_rgb)
    if top_bar:
        gradient_fill(top_bar, T.bg, T.bg2, 0)
        set_solid_alpha(top_bar, 80)
    hl = rect(slide, 0, 0.44, W, 0.06, T.accent_rgb)
    if hl:
        multi_stop_gradient(hl, [(0, T.bg), (40, T.accent), (60, T.accent2), (100, T.bg)], 0)

    # شكراً — hero text
    thanks_y = H * 0.2
    txt(slide, "شكراً جزيلاً", 0.5, thanks_y, W - 1.0, H * 0.28,
        font=_FONT, size=52, bold=True,
        color=T.text_light_rgb, align=PP_ALIGN.CENTER,
        rtl=True, vcenter=True, line_spacing=1.1)

    # "وتقديراً" — secondary
    txt(slide, "وتقديراً", 0.5, thanks_y + H * 0.26, W - 1.0, H * 0.14,
        font=_FONT, size=28, bold=False,
        color=T.accent_rgb, align=PP_ALIGN.CENTER,
        rtl=True, vcenter=True)

    # خط فاصل ذهبي مركزي
    div = rect(slide, W * 0.2, H * 0.56, W * 0.6, 0.055, T.accent_rgb)
    if div:
        multi_stop_gradient(div, [(0, T.bg2), (50, T.accent), (100, T.bg2)], 0)

    # بطاقة معلومات الطالب
    info_y = H * 0.6
    info_h = H - info_y - 0.55
    info_x = W * 0.18
    info_w = W * 0.64

    ic = rrect(slide, info_x, info_y, info_w, info_h, T.card_rgb, radius_pct=12)
    if ic:
        multi_stop_gradient(ic, [(0, T.bg2), (100, T.card)], 135)
        shadow(ic, blur=22, dist=7, alpha=0.5)

    # شريط accent أيمن
    vline(slide, info_x + info_w - 0.2, info_y, info_h, T.accent_rgb, thickness=0.2)

    rows = []
    if req.student_name: rows.append(req.student_name)
    if req.supervisor:   rows.append(f"إشراف: {req.supervisor}")
    if req.year:         rows.append(req.year)

    rh = info_h / max(len(rows), 1)
    for i, row in enumerate(rows):
        fs = 15 if i == 0 else 12
        bold = i == 0
        color = T.text_light_rgb if i == 0 else T.muted_rgb
        txt(slide, row, info_x + 0.3, info_y + i * rh, info_w - 0.5, rh,
            font=_FONT, size=fs, bold=bold,
            color=color, align=PP_ALIGN.CENTER, rtl=True, vcenter=True)

    # شريط سفلي
    bot = rect(slide, 0, H - 0.38, W, 0.38, T.bg_rgb)
    if bot:
        gradient_fill(bot, T.bg2, T.bg, 0)
    hl2 = rect(slide, 0, H - 0.38, W, 0.055, T.accent_rgb)
    if hl2:
        multi_stop_gradient(hl2, [(0, T.bg), (40, T.accent), (60, T.accent2), (100, T.bg)], 0)

    # شعار مذكرتي Pro
    txt(slide, "مذكرتي Pro ✦", 0, H - 0.35, W, 0.35,
        font=_FONT, size=9, bold=False,
        color=T.muted_rgb, align=PP_ALIGN.CENTER, rtl=True, vcenter=True)

    return slide

# ══════════════════════════════════════════════════════════════════════
# PIPELINE ENTRY POINT
# ══════════════════════════════════════════════════════════════════════
def build(prs, req: PresentationRequest, T: Theme, font: str = "Cairo"):
    """نقطة دخول المحرك Ultra"""
    global _FONT
    _FONT = font

    slides_built = []
    sl = req.slides

    # حساب العدد الكلي للشرائح
    total = sum([
        1 if sl.cover else 0,
        1 if sl.intro and (req.intro_overview or req.intro_approach) else 0,
        1 if sl.plan and req.chapters else 0,
        1 if sl.problem else 0,
        1 if sl.objectives and (req.objectives or req.hypotheses) else 0,
        1 if sl.importance and req.importance else 0,
        1 if sl.methodology else 0,
        1 if sl.kpi and req.stats else 0,
        1 if sl.results and req.main_results else 0,
        1 if sl.conclusion else 0,
        1 if sl.recommendations and req.recommendations else 0,
        1 if sl.future and req.future_work else 0,
        1 if sl.references and req.references else 0,
        1 if sl.thankyou else 0,
    ])
    req._total_slides = max(total, 1)

    slide_num = 0

    if sl.cover:
        slides_built.append(make_cover(prs, req, T))

    if sl.intro and (req.intro_overview or req.intro_approach):
        slide_num += 1
        slides_built.append(make_intro(prs, req, T))

    if sl.plan and req.chapters:
        slide_num += 1
        slides_built.append(make_plan(prs, req, T))

    if sl.problem:
        slide_num += 1
        slides_built.append(make_problem(prs, req, T))

    if sl.objectives and (req.objectives or req.hypotheses):
        slide_num += 1
        slides_built.append(make_objectives(prs, req, T))

    if sl.importance and req.importance:
        slide_num += 1
        slides_built.append(make_importance(prs, req, T))

    if sl.methodology:
        slide_num += 1
        slides_built.append(make_methodology(prs, req, T))

    if sl.kpi and req.stats:
        slide_num += 1
        slides_built.append(make_kpi(prs, req, T))

    if sl.results and req.main_results:
        slide_num += 1
        slides_built.append(make_results(prs, req, T))

    if sl.conclusion:
        slide_num += 1
        slides_built.append(make_conclusion(prs, req, T))

    if sl.recommendations and req.recommendations:
        slide_num += 1
        slides_built.append(make_recommendations(prs, req, T))

    if sl.future and req.future_work:
        slide_num += 1
        slides_built.append(make_future(prs, req, T))

    if sl.references and req.references:
        slide_num += 1
        slides_built.append(make_references(prs, req, T))

    if sl.thankyou:
        slides_built.append(make_thankyou(prs, req, T))

    return len(slides_built)

# ══════════════════════════════════════════════════════════════════════
# ALIASES — للتوافق مع pipeline
# ══════════════════════════════════════════════════════════════════════
make_stats = make_kpi
make_final = make_thankyou
